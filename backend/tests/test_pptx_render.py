"""Execução do plano sobre o .pptx do usuário (app/services/pptx_render.py).

Os testes vão do modelo até o arquivo final: constroem um template sintético,
importam (o que dá os `src_shape_id` reais), marcam os slots, planejam e
renderizam. As invariantes da Etapa 6 do plano — nada do modelo sobrando,
nenhum anexo perdido em silêncio, contagem de slides igual à planejada — são
verificadas sobre o arquivo GERADO, não sobre a estrutura interna.
"""
import collections.abc  # noqa: F401  (compat python-pptx 0.6.x + py3.12)
import io
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

from app.services import pptx_render as render
from app.services.deck_plan import ActivityContent, build_plan
from app.services.pptx_import import import_pptx_to_layout


def _png(color=(12, 55, 156), size=(400, 300)) -> bytes:
    buffer = io.BytesIO()
    Image.new("RGB", size, color).save(buffer, "PNG")
    return buffer.getvalue()


@pytest.fixture
def template(tmp_path) -> Path:
    """Capa + um molde de conteúdo com texto, data, tabela, imagem e gráfico."""
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    cover = presentation.slides.add_slide(blank)
    caixa = cover.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.2))
    run = caixa.text_frame.paragraphs[0].add_run()
    run.text = "WEEKLY DO MODELO"
    run.font.size, run.font.bold, run.font.name = Pt(40), True, "Arial"
    selo = cover.shapes.add_textbox(Inches(11), Inches(0.3), Inches(1.6), Inches(0.5))
    s_run = selo.text_frame.paragraphs[0].add_run()
    s_run.text = "W29"
    s_run.font.size, s_run.font.name = Pt(18), "Arial"

    mold = presentation.slides.add_slide(blank)
    titulo = mold.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8), Inches(0.8))
    t_run = titulo.text_frame.paragraphs[0].add_run()
    t_run.text = "ATIVIDADE DA SEMANA PASSADA"
    t_run.font.size, t_run.font.bold, t_run.font.name = Pt(24), True, "Arial"

    data = mold.shapes.add_textbox(Inches(11.5), Inches(0.4), Inches(1.3), Inches(0.4))
    d_run = data.text_frame.paragraphs[0].add_run()
    d_run.text = "19/07"
    d_run.font.size, d_run.font.name = Pt(12), "Arial"

    corpo = mold.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6), Inches(2))
    c_run = corpo.text_frame.paragraphs[0].add_run()
    c_run.text = "Descrição antiga que não pode vazar."
    c_run.font.size, c_run.font.name = Pt(14), "Arial"

    tabela = mold.shapes.add_table(3, 2, Inches(0.5), Inches(3.8), Inches(5), Inches(1.5))
    for linha in range(3):
        for coluna in range(2):
            tabela.table.cell(linha, coluna).text = f"velho {linha}{coluna}"

    imagem = mold.shapes.add_picture(
        io.BytesIO(_png((200, 30, 30))), Inches(7), Inches(1.4), Inches(4), Inches(3)
    )
    assert imagem is not None

    dados = CategoryChartData()
    dados.categories = ["A", "B"]
    dados.add_series("Antiga", (1.0, 2.0))
    mold.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7), Inches(4.6),
        Inches(5), Inches(2.5), dados,
    )

    destino = tmp_path / "modelo.pptx"
    presentation.save(str(destino))
    return destino


@pytest.fixture
def layout(template) -> dict:
    """Modelo importado e com os slots marcados, como o usuário faria na tela."""
    marcado = import_pptx_to_layout(str(template))
    capa, molde = marcado["slides"][0], marcado["slides"][1]

    for element in capa["elements"]:
        texto = (element.get("text") or "").strip()
        element["slot"] = "week_label" if texto == "W29" else "title"
    capa["kind"] = "cover"

    for element in molde["elements"]:
        if element["type"] == "table":
            element["slot"] = "table"
        elif element.get("slot") == "chart":
            element["slot"] = "chart"
        elif element["type"] == "image":
            element["slot"] = "image"
        else:
            texto = (element.get("text") or "").strip()
            if texto == "19/07":
                element["slot"] = "activity_date"
            elif texto.startswith("ATIVIDADE"):
                element["slot"] = "title"
            else:
                element["slot"] = "body"
    return marcado


def _textos(caminho: str) -> str:
    presentation = Presentation(caminho)
    partes = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                partes.append(shape.text)
            if getattr(shape, "has_table", False):
                for linha in shape.table.rows:
                    partes.extend(celula.text for celula in linha.cells)
    return "\n".join(partes)


@pytest.fixture
def foto(tmp_path) -> str:
    caminho = tmp_path / "evidencia.png"
    caminho.write_bytes(_png((30, 200, 90), (600, 200)))
    return str(caminho)


def _render(template, layout, atividades, tmp_path, **kwargs):
    plano = build_plan(
        layout, atividades,
        title=kwargs.get("title", "Weekly W33"),
        subtitle=kwargs.get("subtitle", ""),
        week_label=kwargs.get("week_label", "W33"),
    )
    return plano, render.render_plan(
        template, layout, plano, tmp_path / "saida.pptx"
    )


# ───────────────────────────── invariantes ───────────────────────────────────

def test_gera_arquivo_valido_com_os_slides_do_plano(template, layout, tmp_path):
    atividades = [ActivityContent(title="Auditoria de linha", description="Feita.")]
    plano, resultado = _render(template, layout, atividades, tmp_path)

    assert Path(resultado.path).exists()
    assert resultado.slides == len(plano.slides)
    assert len(Presentation(resultado.path).slides) == len(plano.slides)


def test_nao_sobra_conteudo_do_modelo(template, layout, tmp_path):
    """A causa raiz nº 1 da auditoria: semana e atividade antigas vazando."""
    atividades = [ActivityContent(
        title="Auditoria de linha", description="Concluída sem desvios.",
        date_label="20/07",
    )]
    _, resultado = _render(template, layout, atividades, tmp_path)
    saida = _textos(resultado.path)

    assert "ATIVIDADE DA SEMANA PASSADA" not in saida
    assert "Descrição antiga" not in saida
    assert "19/07" not in saida
    assert "W29" not in saida
    assert "velho" not in saida

    assert "Auditoria de linha" in saida
    assert "Concluída sem desvios." in saida
    assert "20/07" in saida
    assert "W33" in saida


def test_slot_sem_conteudo_e_removido(template, layout, tmp_path):
    """Atividade sem anexo: imagem, tabela e gráfico do modelo saem do slide."""
    atividades = [ActivityContent(title="Reunião", description="Alinhamento.")]
    _, resultado = _render(template, layout, atividades, tmp_path)

    presentation = Presentation(resultado.path)
    conteudo = presentation.slides[1]
    assert not any(getattr(s, "has_table", False) for s in conteudo.shapes)
    assert not any(getattr(s, "has_chart", False) for s in conteudo.shapes)
    assert not any(s.shape_type is not None and "PICTURE" in str(s.shape_type)
                   for s in conteudo.shapes)


def test_preenche_tabela_imagem_e_grafico(template, layout, tmp_path, foto):
    atividades = [ActivityContent(
        title="Inspeção",
        description="Amostras coletadas.",
        tables=[{"columns": ["Item", "Valor"], "rows": [["Peça A", "12"], ["Peça B", "7"]]}],
        images=[foto],
        charts=[{"categories": ["Seg", "Ter"], "series": [{"name": "Falhas", "values": [3, 5]}]}],
    )]
    _, resultado = _render(template, layout, atividades, tmp_path)

    presentation = Presentation(resultado.path)
    conteudo = presentation.slides[1]
    tabela = next(s for s in conteudo.shapes if getattr(s, "has_table", False))
    assert tabela.table.cell(0, 0).text == "Item"
    assert tabela.table.cell(1, 0).text == "Peça A"

    grafico = next(s for s in conteudo.shapes if getattr(s, "has_chart", False))
    assert list(grafico.chart.plots[0].categories) == ["Seg", "Ter"]
    assert "Antiga" not in [s.name for s in grafico.chart.series]

    assert "velho" not in _textos(resultado.path)


def test_anexo_que_nao_cabe_vira_slide_de_continuacao(template, layout, tmp_path, foto):
    """Nada é descartado em silêncio: a 2ª imagem ganha um slide."""
    atividades = [ActivityContent(
        title="Evidências", description="Três fotos.", images=[foto, foto, foto],
    )]
    plano, resultado = _render(template, layout, atividades, tmp_path)

    # capa + 3 slides (o molde tem 1 slot de imagem)
    assert resultado.slides == len(plano.slides) == 4
    saida = _textos(resultado.path)
    assert saida.count("Três fotos.") == 1   # descrição não se repete


def test_texto_grande_reduz_o_corpo_e_avisa_so_se_cortar(template, layout, tmp_path):
    longo = " ".join(["Verificação detalhada do processo produtivo."] * 40)
    atividades = [ActivityContent(title="Relato longo", description=longo)]
    _, resultado = _render(template, layout, atividades, tmp_path)

    presentation = Presentation(resultado.path)
    corpo = next(
        s for s in presentation.slides[1].shapes
        if getattr(s, "has_text_frame", False) and "Verificação" in s.text
    )
    tamanhos = [
        run.font.size.pt
        for paragrafo in corpo.text_frame.paragraphs
        for run in paragrafo.runs
        if run.font.size is not None
    ]
    assert tamanhos and max(tamanhos) < 14   # encolheu em relação ao modelo


def test_texto_impossivel_de_caber_avisa(template, layout, tmp_path):
    gigante = " ".join(["palavra"] * 4000)
    atividades = [ActivityContent(title="Excesso", description=gigante)]
    _, resultado = _render(template, layout, atividades, tmp_path)

    assert any("cortado" in aviso for aviso in resultado.warnings)


def test_avisos_do_plano_chegam_ao_resultado(template, layout, tmp_path):
    """Modelo sem slot para o anexo: o aviso do planejador tem que sobreviver."""
    for slide in layout["slides"]:
        for element in slide["elements"]:
            if element.get("slot") in ("table", "image", "chart"):
                element["slot"] = "static"
    atividades = [ActivityContent(
        title="Com anexo", description="x",
        tables=[{"columns": ["a"], "rows": [["1"]]}],
    )]
    _, resultado = _render(template, layout, atividades, tmp_path)

    assert any("não couberam" in aviso for aviso in resultado.warnings)


def test_slot_static_do_modelo_e_preservado(template, layout, tmp_path):
    for element in layout["slides"][1]["elements"]:
        if (element.get("text") or "").strip().startswith("Descrição antiga"):
            element["slot"] = "static"
    atividades = [ActivityContent(title="Só título", description="ignorada")]
    _, resultado = _render(template, layout, atividades, tmp_path)

    assert "Descrição antiga" in _textos(resultado.path)


def test_modelo_ausente_da_erro_claro(layout, tmp_path):
    plano = build_plan(layout, [ActivityContent(title="x")], title="t")
    with pytest.raises(Exception) as erro:
        render.render_plan(tmp_path / "nao_existe.pptx", layout, plano,
                           tmp_path / "saida.pptx")
    assert "modelo" in str(erro.value).lower()
