"""Conversor .pptx → DeckLayout (modelo da IA a partir de PPT enviado).

Gera um .pptx que imita um weekly feito à mão (capa + slide com título,
descrição, tabela e imagem) e valida a fidelidade da conversão e da clonagem.
"""
import collections.abc  # noqa: F401  (blindagem do pptx antigo no 3.12)
import io
from types import SimpleNamespace

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.services.pptx_import import PptxImportError, import_pptx_to_layout


@pytest.fixture(scope="module")
def modelo_pptx(tmp_path_factory):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # capa
    s1 = prs.slides.add_slide(blank)
    tb = s1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = "Weekly W30 — OQC"
    r.font.size = Pt(40); r.font.bold = True
    r.font.color.rgb = RGBColor(0x0C, 0x37, 0x9C); r.font.name = "Georgia"

    # conteúdo
    s2 = prs.slides.add_slide(blank)
    t = s2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1))
    rt = t.text_frame.paragraphs[0].add_run()
    rt.text = "Auditoria de processo na linha 3"
    rt.font.size = Pt(28); rt.font.bold = True
    d = s2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6), Inches(3))
    d.text_frame.paragraphs[0].add_run().text = "Verificados 12 pontos; 2 NCs."
    s2.shapes.add_table(3, 2, Inches(7.2), Inches(1.8), Inches(5), Inches(2))
    buf = io.BytesIO()
    from PIL import Image
    Image.new("RGB", (400, 250), (30, 90, 180)).save(buf, "PNG")
    buf.seek(0)
    s2.shapes.add_picture(buf, Inches(7.2), Inches(4.2), Inches(5), Inches(2.5))

    path = tmp_path_factory.mktemp("pptx") / "modelo.pptx"
    prs.save(str(path))
    return str(path)


def test_conversion_fidelity(modelo_pptx):
    layout = import_pptx_to_layout(modelo_pptx)
    assert len(layout["slides"]) == 2
    cover, content = layout["slides"]
    assert cover["kind"] == "cover"
    assert content["kind"] == "custom"

    # Capa: título fiel (fonte, tamanho, cor)
    title = cover["elements"][0]
    assert title["type"] == "text"
    assert title["text"] == "Weekly W30 — OQC"
    assert title["font_size"] == 40
    assert title["bold"] is True
    assert title["font_family"] == "Georgia"
    assert title["color"] == "#0C379C"
    # posição em frações (1"/13.333", 2.5"/7.5")
    assert abs(title["x"] - 0.075) < 0.01
    assert abs(title["y"] - 0.3333) < 0.01

    # Conteúdo: título + descrição + slot de tabela + slot de imagem
    types = [e["type"] for e in content["elements"]]
    assert "text" in types
    assert types.count("table") == 1
    assert types.count("image") == 1


def test_clone_from_imported_template(modelo_pptx):
    from app.api.routes.ai_features import _template_deck

    layout = import_pptx_to_layout(modelo_pptx)
    acts = [SimpleNamespace(
        title="Nova atividade", description="Descrição nova.",
        attachments=[
            SimpleNamespace(id="t1", kpi_data={"table": {"columns": ["x"], "rows": []}},
                            file_type="doc", mime_type=""),
            SimpleNamespace(id="i1", kpi_data=None, file_type="image", mime_type="image/png"),
        ],
    )]
    deck = _template_deck(layout, acts, "Weekly W34", "17–23 ago")
    cover, content = deck["slides"][0], deck["slides"][1]
    # capa mantém o estilo do modelo, troca o texto
    ct = cover["elements"][0]
    assert ct["text"] == "Weekly W34"
    assert ct["font_family"] == "Georgia" and ct["font_size"] == 40
    # conteúdo: título trocado + slots preenchidos com os anexos da semana
    att_ids = {e.get("attachment_id") for e in content["elements"] if e.get("attachment_id")}
    assert "t1" in att_ids and "i1" in att_ids


def test_invalid_file_raises():
    import tempfile
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        f.write(b"isto nao e um pptx")
        path = f.name
    with pytest.raises(PptxImportError):
        import_pptx_to_layout(path)


# ── Slots (Etapa 2) ──────────────────────────────────────────────────────────

def _slots_por_texto(slide):
    return {e.get("text"): e.get("slot") for e in slide["elements"] if e["type"] == "text"}


def test_import_suggests_slots_and_anchors_to_the_source_shape(tmp_path):
    """Cada elemento precisa da âncora (slide + shape_id) para o exportador
    por mutação achar o shape certo, e de uma sugestão de papel."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    from app.services.pptx_import import import_pptx_to_layout

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    def caixa(texto, x, y, w, h, size):
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        run = shape.text_frame.paragraphs[0].add_run()
        run.text = texto
        run.font.size = Pt(size)
        return shape

    caixa("W29", 0.4, 0.3, 1.0, 0.5, 28)                  # selo da semana
    caixa("EVOLUÇÃO DOS KPIS", 2.0, 0.3, 8.0, 0.6, 28)    # título de verdade
    caixa("Texto longo da descrição da atividade desta semana, com detalhe.", 0.5, 1.5, 6.0, 3.0, 14)
    caixa("QWI · rodapé", 0.5, 7.0, 3.0, 0.3, 10)
    slide.shapes.add_table(2, 2, Inches(7), Inches(1.5), Inches(4), Inches(1))

    caminho = tmp_path / "modelo.pptx"
    presentation.save(str(caminho))
    layout = import_pptx_to_layout(str(caminho))

    elementos = layout["slides"][0]["elements"]
    assert all(e.get("src_shape_id") for e in elementos)
    assert all(e.get("src_slide") == 0 for e in elementos)

    slots = _slots_por_texto(layout["slides"][0])
    assert slots["W29"] == "week_label"          # não é título, é rótulo de semana
    assert slots["EVOLUÇÃO DOS KPIS"] == "title"  # desempate pela largura
    assert slots["Texto longo da descrição da atividade desta semana, com detalhe."] == "body"
    assert slots["QWI · rodapé"] == "static"      # decoração: repete, nunca recebe conteúdo
    assert [e["slot"] for e in elementos if e["type"] == "table"] == ["table"]


def test_native_chart_gets_its_own_slot(tmp_path):
    """Gráfico nativo não é 'mais uma imagem': o exportador preenche os dados
    dele preservando o estilo do modelo."""
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    from app.services.pptx_import import import_pptx_to_layout

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    dados = CategoryChartData()
    dados.categories = ["A", "B"]
    dados.add_series("S", (1.0, 2.0))
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1),
                           Inches(6), Inches(4), dados)
    shape = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(4), Inches(0.5))
    shape.text_frame.paragraphs[0].add_run().text = "Legenda"

    caminho = tmp_path / "grafico.pptx"
    presentation.save(str(caminho))
    layout = import_pptx_to_layout(str(caminho))

    slots = [e.get("slot") for e in layout["slides"][0]["elements"]]
    assert "chart" in slots


def test_slide_without_a_title_gets_no_title_suggestion(tmp_path):
    """Um palpite errado é pior que palpite nenhum: o usuário pode não notar e
    o slide sai com o conteúdo trocado. Molde de gráfico + parágrafo não tem
    título — e o número da página não pode virar um."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    from app.services.pptx_import import import_pptx_to_layout

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    def caixa(texto, x, y, w, h, size):
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        run = shape.text_frame.paragraphs[0].add_run()
        run.text = texto
        run.font.size = Pt(size)

    caixa("3", 12.5, 7.0, 0.4, 0.3, 16)                    # número da página
    caixa("Análise longa da semana com várias observações relevantes.",
          0.5, 4.0, 8.0, 2.0, 16)                          # parágrafo, embaixo
    caminho = tmp_path / "sem_titulo.pptx"
    presentation.save(str(caminho))

    slots = {e.get("text"): e.get("slot")
             for e in import_pptx_to_layout(str(caminho))["slides"][0]["elements"]}
    assert "title" not in slots.values()
    assert slots["3"] == "static"
    assert slots["Análise longa da semana com várias observações relevantes."] == "body"


def test_week_label_is_detected_in_the_cover_subtitle(tmp_path):
    """'Weekly Report - 19/07/2026 a 19/07/2026' precisa ser trocado, senão a
    capa sai com a data da semana do modelo."""
    from app.services.pptx_import import _is_week_label

    assert _is_week_label("Weekly Report - 19/07/2026 a 19/07/2026")
    assert _is_week_label("W29")
    assert _is_week_label("19/07–19/07")
    # data DENTRO de um parágrafo é conteúdo, não rótulo
    assert not _is_week_label(
        "Durante o período de 01/08 a 07/08 tratamos as falhas de bateria "
        "reportadas pela assistência técnica regional."
    )
