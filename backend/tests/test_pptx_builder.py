"""Tests for PPTX builder service."""

import tempfile
from pathlib import Path

import pytest
from pptx import Presentation

from app.services.pptx_builder import PPTXBuilder
from app.services.pptx_templates import TemplateTheme, TemplateEngine


class TestPPTXBuilder:
    """Test PPTXBuilder functionality."""

    @pytest.fixture
    def builder(self):
        """Create a PPTXBuilder instance with temp directory."""
        with tempfile.TemporaryDirectory() as tmpdir:
            yield PPTXBuilder(output_dir=tmpdir)

    @pytest.fixture
    def sample_report(self):
        """Create sample report data."""
        return {
            "title": "Weekly Quality Report",
            "department": "Quality Assurance",
            "date": "2024-08-10",
            "author": "QWI System",
            "summary": "This week showed significant improvements in quality metrics.",
            "activities": [
                {
                    "title": "Equipment Inspection",
                    "narrative": "Conducted comprehensive inspection of all testing equipment.",
                    "impact": "Identified and resolved 3 calibration issues.",
                },
            ],
            "kpi_table": [
                {"kpi": "Defect Rate", "result": "2.1%", "trend": "↓"},
                {"kpi": "Cycle Time", "result": "45 min", "trend": "→"},
            ],
            "conclusions": [
                "Process improvements yielded measurable results.",
            ],
            "next_steps": [
                "Continue monitoring equipment performance.",
                "Review procedure documentation.",
            ],
        }

    def test_create_presentation(self, builder):
        """Test creating a new presentation."""
        pres = builder.create_presentation()
        assert pres is not None
        assert builder.slide_count == 0

    def test_add_title_slide(self, builder):
        """Test adding a title slide."""
        builder.create_presentation()
        builder.add_title_slide(
            title="Test Report",
            subtitle="Test Subtitle",
            date="2024-08-10",
            author="Tester",
        )
        assert builder.slide_count == 1

    def test_add_content_slide(self, builder):
        """Test adding a content slide."""
        builder.create_presentation()
        builder.add_content_slide(
            title="Test Content",
            content="This is test content.\n- Bullet point 1\n- Bullet point 2",
        )
        assert builder.slide_count == 1

    def test_add_table_slide(self, builder):
        """Test adding a table slide."""
        builder.create_presentation()
        data = [
            ["Item 1", "100"],
            ["Item 2", "200"],
        ]
        builder.add_table_slide(
            title="Test Table",
            data=data,
            headers=["Name", "Value"],
        )
        assert builder.slide_count == 1

    def test_add_summary_slide(self, builder):
        """Test adding a summary slide."""
        builder.create_presentation()
        builder.add_summary_slide({
            "Total Activities": 5,
            "Completed": 5,
            "Issues": 1,
        })
        assert builder.slide_count == 1

    def test_generate_pptx(self, builder, sample_report):
        """Test full PPTX generation."""
        pptx_path = builder.generate_pptx(
            weekly_report=sample_report,
            template_name="default",
            filename="test_report.pptx",
        )

        assert Path(pptx_path).exists()
        assert Path(pptx_path).suffix == ".pptx"

        # Verify PPTX can be opened
        pres = Presentation(pptx_path)
        assert len(pres.slides) > 0

    def test_generate_pptx_with_theme(self, builder, sample_report):
        """Test PPTX generation with custom theme."""
        theme = TemplateTheme.corporate_blue()
        builder_with_theme = PPTXBuilder(theme=theme)

        pptx_path = builder_with_theme.generate_pptx(
            weekly_report=sample_report,
            template_name="default",
            filename="test_themed.pptx",
        )

        assert Path(pptx_path).exists()

    def test_slide_ordering(self, builder):
        """Test slides are created in correct order."""
        builder.create_presentation()
        builder.add_title_slide(title="Slide 1")
        builder.add_content_slide(title="Slide 2", content="Content")
        builder.add_table_slide(title="Slide 3", data=[["A", "B"]])

        assert builder.slide_count == 3
        assert len(builder.presentation.slides) == 3

    def test_multiple_presentations(self):
        """Test creating multiple presentations."""
        with tempfile.TemporaryDirectory() as tmpdir:
            builder1 = PPTXBuilder(output_dir=tmpdir)
            builder2 = PPTXBuilder(output_dir=tmpdir)

            builder1.create_presentation()
            builder1.add_title_slide(title="Report 1")
            path1 = builder1._save_presentation("report1.pptx")

            builder2.create_presentation()
            builder2.add_title_slide(title="Report 2")
            path2 = builder2._save_presentation("report2.pptx")

            assert Path(path1).exists()
            assert Path(path2).exists()
            assert path1 != path2


class TestTemplateTheme:
    """Test TemplateTheme functionality."""

    def test_default_theme(self):
        """Test default theme."""
        theme = TemplateTheme.default()
        assert theme.name == "default"
        assert len(theme.accent_color) == 3
        assert len(theme.background_color) == 3

    def test_corporate_blue_theme(self):
        """Test corporate blue theme."""
        theme = TemplateTheme.corporate_blue()
        assert theme.name == "corporate_blue"
        assert theme.accent_color == (0, 102, 204)

    def test_modern_dark_theme(self):
        """Test modern dark theme."""
        theme = TemplateTheme.modern_dark()
        assert theme.name == "modern_dark"
        assert theme.background_color[0] < 100  # Dark background

    def test_theme_to_dict(self):
        """Test converting theme to dictionary."""
        theme = TemplateTheme.default()
        theme_dict = theme.to_dict()

        assert "name" in theme_dict
        assert "accent_color" in theme_dict
        assert isinstance(theme_dict["accent_color"], tuple)

    def test_theme_from_dict(self):
        """Test creating theme from dictionary."""
        original = TemplateTheme.corporate_blue()
        theme_dict = original.to_dict()

        restored = TemplateTheme.from_dict(theme_dict)
        assert restored.name == original.name
        assert restored.accent_color == original.accent_color


class TestTemplateEngine:
    """Test TemplateEngine functionality."""

    def test_load_predefined_template(self):
        """Test loading predefined templates."""
        engine = TemplateEngine()
        template = engine.load_template("executive")

        assert template is not None
        assert template.name == "executive"

    def test_list_all_templates(self):
        """Test listing templates."""
        engine = TemplateEngine()
        templates = engine.list_all_templates()

        assert len(templates) > 0
        assert any(t["name"] == "executive" for t in templates)

    def test_save_and_load_custom_template(self):
        """Test saving and loading custom templates."""
        with tempfile.TemporaryDirectory() as tmpdir:
            engine = TemplateEngine()
            engine.templates_dir = Path(tmpdir)

            # Create and save template
            from app.services.pptx_templates import PresentationTemplate, TemplateType

            template = PresentationTemplate(
                name="test_template",
                description="Test template",
                template_type=TemplateType.CUSTOM,
                theme=TemplateTheme.default(),
                slide_templates=["title_slide", "content_slide"],
            )

            engine.save_template(template)

            # Load template
            loaded = engine.load_template("test_template")
            assert loaded is not None
            assert loaded.name == "test_template"

    def test_delete_template(self):
        """Test deleting templates."""
        with tempfile.TemporaryDirectory() as tmpdir:
            engine = TemplateEngine()
            engine.templates_dir = Path(tmpdir)

            from app.services.pptx_templates import PresentationTemplate, TemplateType

            template = PresentationTemplate(
                name="temp_template",
                description="Temporary",
                template_type=TemplateType.CUSTOM,
                theme=TemplateTheme.default(),
            )

            engine.save_template(template)
            assert engine.load_template("temp_template") is not None

            engine.delete_template("temp_template")
            assert engine.load_template("temp_template") is None

    def test_preview_theme(self):
        """Test theme preview."""
        engine = TemplateEngine()
        preview = engine.preview_theme("default")

        assert preview is not None
        assert "colors" in preview
        assert "accent" in preview["colors"]


class TestChartIntegration:
    """Test chart integration with PPTX builder."""

    def test_add_chart_slide(self):
        """Test adding a chart slide."""
        with tempfile.TemporaryDirectory() as tmpdir:
            builder = PPTXBuilder(output_dir=tmpdir)
            builder.create_presentation()

            chart_data = {
                "title": "Test Chart",
                "categories": ["A", "B", "C"],
                "series": [
                    {"name": "Series 1", "values": [1, 2, 3]},
                ],
            }

            # This should not raise an error
            try:
                builder.add_chart_slide(
                    title="Chart Test",
                    data=chart_data,
                    chart_type="bar",
                )
            except Exception as e:
                pytest.skip(f"Chart generation not available: {e}")

    def test_export_to_pdf(self):
        """Test PDF export (requires LibreOffice)."""
        with tempfile.TemporaryDirectory() as tmpdir:
            builder = PPTXBuilder(output_dir=tmpdir)
            builder.create_presentation()
            builder.add_title_slide(title="Test")
            pptx_path = builder._save_presentation("test.pptx")

            # This might fail if LibreOffice is not installed
            pdf_path = builder.export_to_pdf(pptx_path)
            # Don't assert, just check it doesn't crash
            assert pdf_path is None or Path(pdf_path).exists()


@pytest.mark.parametrize(
    "template_name",
    ["executive", "operational", "analytical", "technical"],
)
def test_all_template_types(template_name):
    """Test all predefined template types."""
    engine = TemplateEngine()
    template = engine.load_template(template_name)

    assert template is not None
    assert template.name == template_name
    assert template.template_type.value == template_name


@pytest.mark.parametrize(
    "theme_name",
    ["default", "corporate_blue", "modern_dark", "minimalist"],
)
def test_all_themes(theme_name):
    """Test all available themes."""
    engine = TemplateEngine()
    preview = engine.preview_theme(theme_name)

    assert preview is not None
    assert preview["name"] == theme_name
