"""Primitivas de mutação do PPTX (app/services/pptx_mutate.py).

O fixture monta um template sintético COMPLETO (texto formatado, tabela,
imagem e gráfico nativo) para a suíte não depender de nenhum arquivo local.
Quando o .pptx real do usuário existe na máquina, o último teste roda também
contra ele.
"""
import collections.abc  # noqa: F401  (compat python-pptx 0.6.x + py3.12)
import io
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

from app.services import pptx_mutate as mut


def _png(color=(12, 55, 156), size=(400, 300)) -> bytes:
    buffer = io.BytesIO()
    Image.new("RGB", size, color).save(buffer, "PNG")
    return buffer.getvalue()


@pytest.fixture
def template(tmp_path) -> Path:
    """Template sintético com um slide de capa e um slide-molde completo."""
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    cover = presentation.slides.add_slide(blank)
    box = cover.shapes.add_textbox(Inches(1), Inches(3), Inches(10), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "TÍTULO DO MODELO"
    run.font.size, run.font.bold = Pt(40), True
    run.font.color.rgb = RGBColor(0x0C, 0x37, 0x9C)
    run.font.name = "Georgia"

    mold = presentation.slides.add_slide(blank)
    title = mold.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8), Inches(0.8))
    t_run = title.text_frame.paragraphs[0].add_run()
    t_run.text = "TÍTULO DA ATIVIDADE"
    t_run.font.size, t_run.font.bold = Pt(24), True
    t_run.font.color.rgb = RGBColor(0x0C, 0x37, 0x9C)

    body = mold.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6), Inches(2))
    b_run = body.text_frame.paragraphs[0].add_run()
    b_run.text = "Descrição da semana passada."
    b_run.font.size = Pt(14)
    b_run.font.color.rgb = RGBColor(0x1F, 0x29, 0x37)

    table = mold.shapes.add_table(3, 2, Inches(0.5), Inches(3.6), Inches(5), Inches(1.5))
    for r in range(3):
        for c in range(2):
            table.table.cell(r, c).text = f"velho {r}{c}"

    mold.shapes.add_picture(io.BytesIO(_png()), Inches(7), Inches(1.4),
                            width=Inches(4), height=Inches(2))

    data = CategoryChartData()
    data.categories = ["A", "B"]
    data.add_series("Modelo", (1.0, 2.0))
    mold.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7), Inches(3.6),
                          Inches(5), Inches(2.5), data)

    mold.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(0.2),
                          Inches(0.1), Inches(6))

    path = tmp_path / "modelo.pptx"
    presentation.save(str(path))
    return path


def _roundtrip(presentation, tmp_path, name="saida.pptx"):
    """Salva e reabre — garante que o arquivo é válido de verdade."""
    out = tmp_path / name
    mut.save(presentation, out)
    return Presentation(str(out))


def _kinds(slide) -> dict:
    counts = {"texto": 0, "imagem": 0, "tabela": 0, "grafico": 0, "outro": 0}
    for shape in slide.shapes:
        if getattr(shape, "has_chart", False):
            counts["grafico"] += 1
        elif getattr(shape, "has_table", False):
            counts["tabela"] += 1
        elif shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE):
            counts["imagem"] += 1
        elif getattr(shape, "has_text_frame", False) and shape.text.strip():
            counts["texto"] += 1
        else:
            counts["outro"] += 1
    return counts


def _first(slide, predicate):
    return next(s for s in slide.shapes if predicate(s))


def test_duplicate_slide_keeps_every_kind_of_shape(template, tmp_path):
    presentation = mut.open_presentation(template)
    mut.duplicate_slide(presentation, 1)
    reopened = _roundtrip(presentation, tmp_path)

    assert mut.slide_count(reopened) == 3
    assert _kinds(reopened.slides[1]) == _kinds(reopened.slides[2])
    assert _kinds(reopened.slides[2])["grafico"] == 1


def test_duplicated_image_points_to_a_valid_blob(template, tmp_path):
    """Se os rIds não forem remapeados, a imagem do clone aponta para o
    lugar errado — falha silenciosa que só aparece ao abrir o arquivo."""
    presentation = mut.open_presentation(template)
    mut.duplicate_slide(presentation, 1)
    reopened = _roundtrip(presentation, tmp_path)

    picture = _first(reopened.slides[2], lambda s: s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    assert len(picture.image.blob) > 0
    assert picture.image.ext == "png"


def test_chart_data_of_the_copy_does_not_leak_into_the_original(template, tmp_path):
    """O slide duplicado compartilha a part do gráfico: sem clonar, escrever
    nos dados de um alteraria os dois."""
    presentation = mut.open_presentation(template)
    copy_slide = mut.duplicate_slide(presentation, 1)
    frame = _first(copy_slide, lambda s: getattr(s, "has_chart", False))
    mut.set_chart_data(frame, copy_slide, ["Jan", "Fev", "Mar"], [("Semana", (10, 20, 30))])

    reopened = _roundtrip(presentation, tmp_path)
    original = _first(reopened.slides[1], lambda s: getattr(s, "has_chart", False)).chart
    clone = _first(reopened.slides[2], lambda s: getattr(s, "has_chart", False)).chart

    assert [str(c) for c in original.plots[0].categories] == ["A", "B"]
    assert [str(c) for c in clone.plots[0].categories] == ["Jan", "Fev", "Mar"]
    assert list(clone.plots[0].series[0].values) == [10.0, 20.0, 30.0]
    assert clone.chart_type == original.chart_type   # estilo do modelo preservado


def test_set_text_preserves_formatting(template, tmp_path):
    presentation = mut.open_presentation(template)
    shape = _first(presentation.slides[1],
                   lambda s: getattr(s, "has_text_frame", False) and "TÍTULO" in s.text)
    mut.set_text(shape, "Análise de falhas na linha 3")

    reopened = _roundtrip(presentation, tmp_path)
    saved = _first(reopened.slides[1],
                   lambda s: getattr(s, "has_text_frame", False) and "Análise" in s.text)
    run = saved.text_frame.paragraphs[0].runs[0]
    assert saved.text == "Análise de falhas na linha 3"
    assert run.font.size == Pt(24)
    assert run.font.bold is True
    assert run.font.color.rgb == RGBColor(0x0C, 0x37, 0x9C)


def test_set_paragraphs_repeats_the_style_of_the_first(template, tmp_path):
    presentation = mut.open_presentation(template)
    shape = _first(presentation.slides[1],
                   lambda s: getattr(s, "has_text_frame", False) and "Descrição" in s.text)
    mut.set_paragraphs(shape, ["Primeira linha", "Segunda linha", "Terceira"], "• ")

    reopened = _roundtrip(presentation, tmp_path)
    saved = _first(reopened.slides[1],
                   lambda s: getattr(s, "has_text_frame", False) and "Primeira" in s.text)
    paragraphs = saved.text_frame.paragraphs
    assert len(paragraphs) == 3
    assert paragraphs[2].runs[0].text == "• Terceira"
    assert all(p.runs[0].font.size == Pt(14) for p in paragraphs)


def test_clear_text_keeps_the_box(template, tmp_path):
    """Slot sem conteúdo na semana: a caixa continua lá (para não desmontar o
    layout), mas o texto do modelo NÃO pode sobrar."""
    presentation = mut.open_presentation(template)
    antes = len(presentation.slides[1].shapes._spTree)
    shape = _first(presentation.slides[1],
                   lambda s: getattr(s, "has_text_frame", False) and "Descrição" in s.text)
    mut.clear_text(shape)

    reopened = _roundtrip(presentation, tmp_path)
    slide = reopened.slides[1]
    assert len(slide.shapes._spTree) == antes
    textos = [s.text for s in slide.shapes if getattr(s, "has_text_frame", False)]
    assert all("Descrição da semana passada" not in t for t in textos)
    assert "" in textos


def test_remove_shape(template, tmp_path):
    presentation = mut.open_presentation(template)
    shape = _first(presentation.slides[1], lambda s: getattr(s, "has_table", False))
    mut.remove_shape(shape)

    reopened = _roundtrip(presentation, tmp_path)
    assert _kinds(reopened.slides[1])["tabela"] == 0
    assert _kinds(reopened.slides[1])["grafico"] == 1   # o resto fica intacto


def test_swap_image_contains_inside_the_original_frame(template, tmp_path):
    presentation = mut.open_presentation(template)
    picture = _first(presentation.slides[1], lambda s: s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    frame = (picture.left, picture.top, picture.width, picture.height)
    mut.swap_image(picture, _png(color=(200, 0, 0), size=(300, 900)))  # bem vertical

    reopened = _roundtrip(presentation, tmp_path)
    saved = _first(reopened.slides[1], lambda s: s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    # cabe dentro da moldura original e continua centralizado
    assert saved.width <= frame[2] and saved.height <= frame[3]
    assert saved.left >= frame[0] and saved.top >= frame[1]
    assert abs((saved.left + saved.width / 2) - (frame[0] + frame[2] / 2)) < 2000
    assert abs(saved.width / saved.height - 300 / 900) < 0.01


def test_fill_table_grows_and_shrinks_rows(template, tmp_path):
    presentation = mut.open_presentation(template)
    frame = _first(presentation.slides[1], lambda s: getattr(s, "has_table", False))
    report = mut.fill_table(frame, ["Sintoma", "Qtd"],
                            [["Tela", "88"], ["Não liga", "66"], ["Bateria", "59"],
                             ["Câmera", "31"]])

    reopened = _roundtrip(presentation, tmp_path)
    table = _first(reopened.slides[1], lambda s: getattr(s, "has_table", False)).table
    assert len(table.rows) == 5                      # cabeçalho + 4
    assert table.cell(0, 0).text == "Sintoma"
    assert table.cell(4, 1).text == "31"
    assert report == {"dropped_rows": 0, "dropped_columns": 0}


def test_fill_table_reports_what_did_not_fit(template, tmp_path):
    presentation = mut.open_presentation(template)
    frame = _first(presentation.slides[1], lambda s: getattr(s, "has_table", False))
    report = mut.fill_table(frame, ["A", "B", "C", "D"], [["1", "2", "3", "4"]])
    assert report["dropped_columns"] == 2            # a tabela do modelo tem 2 colunas


def test_delete_and_move_slide(template, tmp_path):
    presentation = mut.open_presentation(template)
    mut.duplicate_slide(presentation, 1)
    mut.move_slide(presentation, 2, 0)
    mut.delete_slide(presentation, 2)

    reopened = _roundtrip(presentation, tmp_path)
    assert mut.slide_count(reopened) == 2
    assert _kinds(reopened.slides[0])["grafico"] == 1  # o clone foi para a frente


def test_open_presentation_gives_a_clear_error(tmp_path):
    quebrado = tmp_path / "nao_e_pptx.pptx"
    quebrado.write_bytes(b"isto nao e um pptx")
    with pytest.raises(mut.PptxMutateError) as error:
        mut.open_presentation(quebrado)
    assert "modelo" in str(error.value).lower()


REAL = Path("uploads/pptx_templates/b0c873a3-2786-4934-8709-a96ef2e2d01c/None.pptx")


@pytest.mark.skipif(not REAL.exists(), reason="template real não está nesta máquina")
def test_duplication_on_the_real_user_template(tmp_path):
    presentation = mut.open_presentation(REAL)
    antes = mut.slide_count(presentation)
    mut.duplicate_slide(presentation, 1)
    reopened = _roundtrip(presentation, tmp_path, "real.pptx")
    assert mut.slide_count(reopened) == antes + 1
    assert _kinds(reopened.slides[1]) == _kinds(reopened.slides[antes])
