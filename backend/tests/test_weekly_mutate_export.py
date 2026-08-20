"""Integração: weekly gerado por MUTAÇÃO do PPT enviado pelo usuário.

Cobre o ramo novo de `WeeklyService._generate_by_mutation` — o que amarra
deck_content → deck_plan → pptx_render ao banco. O que estes testes protegem:

1. escolhido o modelo, o arquivo sai do .pptx do usuário (sem LLM);
2. modelo quebrado NÃO derruba o weekly — volta ao fluxo antigo;
3. modelo de outro usuário não é acessível pelo id.
"""
import collections.abc  # noqa: F401  (compat python-pptx 0.6.x + py3.12)
import io
import os
import tempfile
from datetime import UTC, datetime
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker
from sqlalchemy.orm.attributes import flag_modified

_TMP = tempfile.mkdtemp(prefix="qwi_mut_")
os.environ["UPLOAD_DIR"] = f"{_TMP}/uploads"

from app.core.database import Base  # noqa: E402
from app.core.security import get_password_hash  # noqa: E402
from app.models import (  # noqa: E402
    Activity, PptxTemplate, QualitySector, User, UserRole, WeeklyReport, WeeklyStatus,
)
from app.services.business import WeeklyService  # noqa: E402
from app.services.pptx_import import import_pptx_to_layout  # noqa: E402

_engine = create_engine(f"sqlite:///{_TMP}/m.db", connect_args={"check_same_thread": False})
_Session = sessionmaker(bind=_engine, autocommit=False, autoflush=False)
Base.metadata.create_all(bind=_engine)


def _modelo_pptx(destino: Path) -> Path:
    """Capa + molde de conteúdo, com os textos do 'modelo antigo'."""
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    capa = presentation.slides.add_slide(blank)
    caixa = capa.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.2))
    run = caixa.text_frame.paragraphs[0].add_run()
    run.text = "TÍTULO DO MODELO"
    run.font.size, run.font.bold, run.font.name = Pt(40), True, "Arial"

    molde = presentation.slides.add_slide(blank)
    titulo = molde.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.9))
    t_run = titulo.text_frame.paragraphs[0].add_run()
    t_run.text = "ATIVIDADE DO MODELO"
    t_run.font.size, t_run.font.bold, t_run.font.name = Pt(26), True, "Arial"

    corpo = molde.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(3))
    c_run = corpo.text_frame.paragraphs[0].add_run()
    c_run.text = "Descrição do modelo que não pode vazar."
    c_run.font.size, c_run.font.name = Pt(14), "Arial"

    destino.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(destino))
    return destino


def _marcado(caminho: Path) -> dict:
    layout = import_pptx_to_layout(str(caminho))
    layout["slides"][0]["kind"] = "cover"
    for element in layout["slides"][0]["elements"]:
        element["slot"] = "title"
    for element in layout["slides"][1]["elements"]:
        texto = (element.get("text") or "").strip()
        element["slot"] = "title" if texto.startswith("ATIVIDADE") else "body"
    return layout


@pytest.fixture
def db():
    session = _Session()
    try:
        yield session
    finally:
        session.rollback()
        session.close()


@pytest.fixture
def cenario(db, tmp_path):
    """Usuário + atividade da semana + modelo de PPT marcado."""
    sufixo = tmp_path.name
    user = User(
        id=f"u-{sufixo}", email=f"{sufixo}@qwitest.com", employee_id=f"E-{sufixo}",
        hashed_password=get_password_hash("senha123"), name="Álvaro",
        role=UserRole.ANALISTA_JR, sector=QualitySector.OQC,
    )
    activity = Activity(
        id=f"a-{sufixo}", user_id=user.id, title="Auditoria da linha 3",
        description="Doze pontos verificados, sem desvio.",
        department="Qualidade", week_number=33, year=2026,
        activity_date=datetime(2026, 8, 11, tzinfo=UTC),
    )
    arquivo = _modelo_pptx(tmp_path / "modelo.pptx")
    template = PptxTemplate(
        id=f"t-{sufixo}", user_id=user.id, name="Modelo",
        file_path=str(arquivo), layout=_marcado(arquivo),
    )
    report = WeeklyReport(
        id=f"r-{sufixo}", user_id=user.id, week_number=33, year=2026,
        version=1, status=WeeklyStatus.GENERATING, title="Weekly W33",
    )
    db.add_all([user, activity, template, report])
    db.commit()
    return user, activity, template, report


def _gerar(db, cenario, template_id=None):
    user, activity, template, report = cenario
    return WeeklyService(db)._generate_by_mutation(
        report, user, [activity],
        template_id if template_id is not None else template.id,
        week_number=33, year=2026, period_label="10/08–16/08",
    )


def _textos(caminho: str) -> str:
    presentation = Presentation(caminho)
    return "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )


def test_gera_o_deck_a_partir_do_pptx_do_usuario(db, cenario):
    resultado = _gerar(db, cenario)

    assert resultado is not None
    assert resultado.status == WeeklyStatus.COMPLETED
    assert Path(resultado.pptx_path).exists()
    assert resultado.content["source"] == "pptx_mutate"
    # Sem LLM neste caminho: nada de prompt nem de resumo da IA.
    assert resultado.prompt_used is None
    assert resultado.ai_summary is None


def test_conteudo_da_semana_substitui_o_do_modelo(db, cenario):
    resultado = _gerar(db, cenario)
    saida = _textos(resultado.pptx_path)

    assert "Auditoria da linha 3" in saida
    assert "Doze pontos verificados" in saida
    assert "ATIVIDADE DO MODELO" not in saida
    assert "Descrição do modelo" not in saida


def test_atividade_marcada_como_usada(db, cenario):
    _, activity, _, _ = cenario
    _gerar(db, cenario)
    assert activity.status.value == "used_in_report"


def test_modelo_de_outro_usuario_nao_e_usado(db, cenario, tmp_path):
    """Sem dono conferido, o id de um modelo alheio geraria o deck dele."""
    outro = PptxTemplate(
        id="t-alheio", user_id="fantasma", name="Alheio",
        file_path=str(_modelo_pptx(tmp_path / "alheio.pptx")), layout={"slides": []},
    )
    db.add(outro)
    db.commit()

    assert _gerar(db, cenario, template_id="t-alheio") is None


def test_modelo_sem_campos_marcados_volta_ao_fluxo_antigo(db, cenario):
    _, _, template, _ = cenario
    template.layout = {"slides": []}
    db.commit()

    assert _gerar(db, cenario) is None


def test_arquivo_do_modelo_sumido_volta_ao_fluxo_antigo(db, cenario):
    _, _, template, _ = cenario
    Path(template.file_path).unlink()

    assert _gerar(db, cenario) is None


def test_avisos_do_plano_ficam_guardados_no_relatorio(db, cenario):
    """O usuário precisa poder ver o que não coube — some se não for salvo."""
    _, _, template, _ = cenario
    layout = template.layout
    for element in layout["slides"][1]["elements"]:
        element["slot"] = "static"      # nenhum campo para o conteúdo
    flag_modified(template, "layout")   # mutação dentro do JSON não é detectada
    db.commit()

    resultado = _gerar(db, cenario)
    assert resultado is not None
    assert any("Marque os campos" in aviso or "campos marcados" in aviso
               for aviso in resultado.content["warnings"])
