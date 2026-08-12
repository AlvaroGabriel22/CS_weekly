from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.pptx_service import PptxService

FIELD_DATA = {
    "title": "Weekly Report FIELD",
    "department": "Quality",
    "sector": "FIELD",
    "period_label": "06/07–12/07",
    "summary": "Análise de falha de campo no modelo A20 com medições fora de spec.",
    "highlights": ["49 ocorrências de Missing Screw"],
    "conclusions": ["Causa raiz pendente"],
    "next_steps": ["Auditoria fornecedor W33"],
    "presentation_plan": {
        "layout_profile": "field_case",
        "sidebar": ["synthesis", "highlights", "conclusions", "next_steps"],
        "global_blocks": [
            {
                "type": "chart",
                "title": "Defeitos por tipo",
                "chart_type": "column",
                "categories": ["Missing Screw", "Cosmetic"],
                "series": [{"name": "Qty", "values": [30, 19]}],
                "insight": "Missing Screw concentra 61% dos defeitos.",
            }
        ],
    },
    "activities": [
        {
            "source": 1,
            "title": "FIELD: A20 failure analysis",
            "date": "08/07",
            "narrative": "Material fora de escopo identificado na auditoria de lote.",
            "impact": "Hold de produção até lote melhorado W33.",
            "facts": ["49 ocorrências totais", "Lote fora de escopo"],
            "actions": ["Solicitar lote melhorado ao fornecedor"],
            "blocks": [
                {
                    "type": "device_info",
                    "fields": {
                        "Modelo": "A20",
                        "Código": "SN-8842",
                        "Falha": "Missing Screw",
                    },
                },
                {
                    "type": "measurement_table",
                    "title": "Medições elétricas",
                    "columns": ["Parâmetro", "Valor", "Unidade", "Limite", "Status"],
                    "rows": [
                        ["Corrente standby", "12", "mA", "15", "OK"],
                        ["Torque parafuso", "0.32", "Nm", "0.45", "NOK"],
                    ],
                },
                {
                    "type": "countermeasure_table",
                    "title": "Contramedidas",
                    "rows": [
                        {
                            "action": "Hold produção lote atual",
                            "owner": "IQC",
                            "status": "Em andamento",
                            "due": "W29",
                        }
                    ],
                },
            ],
        }
    ],
}


def test_field_report_renders_tables_and_chart(tmp_path: Path):
    image_path = tmp_path / "evidence.png"
    Image.new("RGB", (300, 200), color=(40, 90, 180)).save(image_path)
    data = {
        **FIELD_DATA,
        "activities": [
            {
                **FIELD_DATA["activities"][0],
                "images": [{"path": str(image_path), "caption": "Evidência"}],
            }
        ],
    }

    service = PptxService(output_dir=str(tmp_path))
    output = service.generate(
        "field-test",
        None,
        {},
        data,
        [],
        28,
        2026,
        "Field Analyst",
        "pt",
    )

    presentation = Presentation(output)
    assert len(presentation.slides) >= 2

    tables = [
        shape
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_table
    ]
    assert len(tables) >= 2

    charts = [
        shape
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_chart
    ]
    assert len(charts) >= 1
    assert charts[0].chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED

    pictures = [
        shape
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    assert len(pictures) >= 1

    text = " ".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    table_text = " ".join(
        cell.text
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_table
        for row in shape.table.rows
        for cell in row.cells
    )
    combined = f"{text} {table_text}"
    assert "A20" in combined
    assert "Torque" in combined or "Corrente" in combined


def test_transcribe_mode_allows_longer_narrative(tmp_path: Path):
    long_text = " ".join(["Conclusão direta registrada pelo analista."] * 20)
    data = {
        "title": "Weekly",
        "department": "Quality",
        "period_label": "06/07–12/07",
        "activities": [
            {
                "source": 1,
                "title": "Auditoria",
                "date": "08/07",
                "content_mode": "transcribe",
                "narrative": long_text,
                "facts": [],
            }
        ],
    }
    service = PptxService(output_dir=str(tmp_path))
    output = service.generate("transcribe-test", None, {}, data, [], 28, 2026, "Analyst", "pt")
    presentation = Presentation(output)
    text = " ".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "Conclusão direta registrada pelo analista." in text
    assert len(text) > 200


def test_compress_mode_truncates_long_narrative(tmp_path: Path):
    long_text = " ".join(["Frase executiva curta."] * 80)
    data = {
        "title": "Weekly",
        "department": "Quality",
        "period_label": "06/07–12/07",
        "activities": [
            {
                "source": 1,
                "title": "Auditoria",
                "date": "08/07",
                "content_mode": "compress",
                "narrative": long_text,
                "facts": [],
            }
        ],
    }
    service = PptxService(output_dir=str(tmp_path))
    output = service.generate("compress-test", None, {}, data, [], 28, 2026, "Analyst", "pt")
    presentation = Presentation(output)
    text = " ".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert text.count("Frase executiva curta.") < 80

