"""Formatação por trecho e espaçamento entre linhas no renderizador de layout.

O editor deixa o usuário formatar SÓ uma parte do texto (negrito, corpo maior,
outra fonte). Isso chega ao backend em `runs`. Estes testes garantem que os
trechos viram runs de verdade no .pptx — e que quem não usa `runs` continua
saindo exatamente como antes.
"""
import collections.abc  # noqa: F401  (compat python-pptx 0.6.x + py3.12)
from pathlib import Path

import pytest
from pptx import Presentation

from app.services.pptx_layout import PptxLayoutRenderer


def _texto(x=0.05, y=0.1, w=0.8, h=0.2, **extra) -> dict:
    return {"id": "t1", "type": "text", "x": x, "y": y, "w": w, "h": h,
            "font_size": 14, "color": "#1F2937", **extra}


def _gerar(tmp_path: Path, element: dict):
    deck = {"slides": [{"id": "s1", "kind": "custom", "elements": [element]}]}
    caminho = PptxLayoutRenderer(tmp_path).generate("rel-1", deck, {}, {})
    slide = Presentation(caminho).slides[0]
    return next(s for s in slide.shapes if getattr(s, "has_text_frame", False))


def test_sem_runs_mantem_um_formato_para_a_caixa(tmp_path):
    caixa = _gerar(tmp_path, _texto(text="Uma linha só", bold=True, font_size=20))
    runs = caixa.text_frame.paragraphs[0].runs
    assert len(runs) == 1
    assert runs[0].text == "Uma linha só"
    assert runs[0].font.bold is True
    assert runs[0].font.size.pt == 20


def test_runs_viram_trechos_com_formatos_diferentes(tmp_path):
    element = _texto(
        text="Resultado: aprovado",
        runs=[
            {"text": "Resultado: "},
            {"text": "aprovado", "bold": True, "font_size": 22, "color": "#16A34A"},
        ],
    )
    runs = _gerar(tmp_path, element).text_frame.paragraphs[0].runs
    assert [r.text for r in runs] == ["Resultado: ", "aprovado"]
    assert runs[0].font.bold in (False, None)
    assert runs[1].font.bold is True
    assert runs[1].font.size.pt == 22
    assert str(runs[1].font.color.rgb) == "16A34A"


def test_run_sem_formato_proprio_herda_o_da_caixa(tmp_path):
    element = _texto(text="herdado", font_family="Georgia", font_size=18,
                     runs=[{"text": "herdado"}])
    run = _gerar(tmp_path, element).text_frame.paragraphs[0].runs[0]
    assert run.font.name == "Georgia"
    assert run.font.size.pt == 18


def test_quebra_de_linha_dentro_do_run_vira_paragrafo(tmp_path):
    element = _texto(text="linha 1\nlinha 2",
                     runs=[{"text": "linha 1\nlinha 2", "bold": True}])
    frame = _gerar(tmp_path, element).text_frame
    assert [p.text for p in frame.paragraphs] == ["linha 1", "linha 2"]
    assert all(r.font.bold for p in frame.paragraphs for r in p.runs)


def test_linha_em_branco_do_usuario_sobrevive(tmp_path):
    """Parágrafo vazio colapsaria e o espaçamento que o usuário deixou sumiria."""
    element = _texto(text="a\n\nb", runs=[{"text": "a\n\nb"}])
    frame = _gerar(tmp_path, element).text_frame
    assert [p.text for p in frame.paragraphs] == ["a", "", "b"]


def test_espacamento_entre_linhas_e_aplicado(tmp_path):
    caixa = _gerar(tmp_path, _texto(text="a\nb", line_spacing=1.5))
    assert all(p.line_spacing == 1.5 for p in caixa.text_frame.paragraphs)


def test_espacamento_ausente_nao_e_escrito(tmp_path):
    """Sem escolha do usuário, herda do PowerPoint em vez de fixar um valor."""
    caixa = _gerar(tmp_path, _texto(text="a"))
    assert caixa.text_frame.paragraphs[0].line_spacing is None


@pytest.mark.parametrize("valor", [0, -1, 9, "muito", None])
def test_espacamento_invalido_e_ignorado(tmp_path, valor):
    caixa = _gerar(tmp_path, _texto(text="a", line_spacing=valor))
    assert caixa.text_frame.paragraphs[0].line_spacing is None


def test_espacamento_vale_tambem_com_runs(tmp_path):
    element = _texto(text="a\nb", line_spacing=2.0, runs=[{"text": "a\nb"}])
    frame = _gerar(tmp_path, element).text_frame
    assert all(p.line_spacing == 2.0 for p in frame.paragraphs)


def test_runs_vazio_cai_no_caminho_antigo(tmp_path):
    caixa = _gerar(tmp_path, _texto(text="sem trechos", runs=[]))
    assert caixa.text_frame.paragraphs[0].runs[0].text == "sem trechos"


# ── fonte coreana ────────────────────────────────────────────────────────────

def _ea_typeface(run) -> str | None:
    from pptx.oxml.ns import qn
    ea = run.font._rPr.find(qn("a:ea"))
    return ea.get("typeface") if ea is not None else None


def test_malgun_gothic_vai_tambem_para_a_tipografia_east_asian(tmp_path):
    """Sem `ea`, o PowerPoint escreve o coreano com a fonte do TEMA."""
    caixa = _gerar(tmp_path, _texto(text="품질 주간 보고", font_family="Malgun Gothic"))
    run = caixa.text_frame.paragraphs[0].runs[0]
    assert run.font.name == "Malgun Gothic"
    assert _ea_typeface(run) == "Malgun Gothic"


def test_fonte_latina_nao_e_forcada_no_east_asian(tmp_path):
    """Forçar Arial no `ea` viraria caixinhas: o tema resolve melhor."""
    caixa = _gerar(tmp_path, _texto(text="Resultado", font_family="Arial"))
    run = caixa.text_frame.paragraphs[0].runs[0]
    assert run.font.name == "Arial"
    assert _ea_typeface(run) is None


def test_trecho_em_malgun_dentro_de_caixa_latina(tmp_path):
    element = _texto(
        text="Título 제목", font_family="Arial",
        runs=[{"text": "Título "}, {"text": "제목", "font_family": "Malgun Gothic"}],
    )
    runs = _gerar(tmp_path, element).text_frame.paragraphs[0].runs
    assert _ea_typeface(runs[0]) is None
    assert _ea_typeface(runs[1]) == "Malgun Gothic"


def test_ea_fica_logo_depois_de_latin(tmp_path):
    """A ordem do schema é ... latin, ea, cs ...

    Fora dela o PowerPoint recusa o arquivo — e o python-pptx, que é
    tolerante, não acusaria. Por isso conferimos a posição, não só o valor.
    """
    from pptx.oxml.ns import qn

    caixa = _gerar(tmp_path, _texto(text="제목", font_family="Malgun Gothic", bold=True))
    rPr = caixa.text_frame.paragraphs[0].runs[0].font._rPr
    filhos = [child.tag for child in rPr]
    assert filhos.index(qn("a:ea")) == filhos.index(qn("a:latin")) + 1
