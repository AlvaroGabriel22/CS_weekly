"""Upload de modelo de PPT (aba Templates) — higiene de arquivo.

Regressão do bug em que TODO modelo era salvo como "None.pptx": o default da
coluna `id` só é aplicado no INSERT, então o caminho era montado com None. Dois
modelos do mesmo usuário viravam o mesmo arquivo — um sobrescrevia o outro, e
apagar um levava junto o arquivo do outro.
"""
import asyncio
import collections.abc  # noqa: F401  (compat python-pptx 0.6.x + py3.12)
import io
import os
import tempfile
from pathlib import Path

import httpx
import pytest
from pptx import Presentation
from pptx.util import Inches, Pt
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker

_TMP = tempfile.mkdtemp(prefix="qwi_tpl_test_")
os.environ["UPLOAD_DIR"] = f"{_TMP}/uploads"

from app.core.database import Base, get_db  # noqa: E402
from app.main import create_app  # noqa: E402

_test_engine = create_engine(
    f"sqlite:///{_TMP}/tpl.db", connect_args={"check_same_thread": False}
)
_TestSession = sessionmaker(bind=_test_engine, autocommit=False, autoflush=False)
Base.metadata.create_all(bind=_test_engine)

_APP = create_app()


def _override_get_db():
    db = _TestSession()
    try:
        yield db
    finally:
        db.close()


_APP.dependency_overrides[get_db] = _override_get_db


class SyncClient:
    def __init__(self):
        self._c = httpx.AsyncClient(
            transport=httpx.ASGITransport(app=_APP),
            base_url="http://test", timeout=30,
        )

    def _run(self, coro):
        return asyncio.get_event_loop().run_until_complete(coro)

    def get(self, url, **kw):
        return self._run(self._c.get(url, **kw))

    def post(self, url, **kw):
        return self._run(self._c.post(url, **kw))

    def patch(self, url, **kw):
        return self._run(self._c.patch(url, **kw))

    def delete(self, url, **kw):
        return self._run(self._c.delete(url, **kw))


_UPLOADS = Path(_TMP) / "uploads"


@pytest.fixture(autouse=True)
def _isolate_uploads(monkeypatch):
    """Aponta o UPLOAD_DIR da rota para um diretório temporário.

    Só mexer em `os.environ` não basta: o `settings` é um singleton já criado
    quando o conftest importa o app, e a rota guarda a referência no import.
    Sem isso, o teste gravaria dentro de `uploads/` de desenvolvimento.
    """
    from app.api.routes import pptx_templates as rota

    monkeypatch.setattr(rota.settings, "UPLOAD_DIR", str(_UPLOADS), raising=False)
    yield


def _templates_dir() -> Path:
    return _UPLOADS / "pptx_templates"


def _token(c, sufixo: str) -> str:
    """Um usuário por teste — o limite é de 2 modelos por usuário."""
    email = f"tpl{sufixo}@empresa.com"
    c.post("/api/auth/register", json={
        "name": f"Dono do Modelo {sufixo}", "email": email, "password": "senha123",
        "password_confirm": "senha123", "employee_id": f"TPL{sufixo}",
        "role": "Analista Jr", "sector": "OQC",
    })
    r = c.post("/api/auth/login", json={"email": email, "password": "senha123"})
    return r.json()["access_token"]


def _pptx_bytes(titulo: str) -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = titulo
    run.font.size = Pt(32)
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _upload(c, token, nome: str):
    return c.post(
        "/api/pptx-templates",
        headers={"Authorization": f"Bearer {token}"},
        files={"file": (f"{nome}.pptx", _pptx_bytes(nome.upper()),
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )


def test_two_templates_do_not_share_the_same_file():
    c = SyncClient()
    token = _token(c, "A")

    primeiro = _upload(c, token, "modelo_a")
    segundo = _upload(c, token, "modelo_b")
    assert primeiro.status_code == 201, primeiro.text
    assert segundo.status_code == 201, segundo.text

    ids = {primeiro.json()["id"], segundo.json()["id"]}
    arquivos = sorted(_templates_dir().rglob("*.pptx"))
    meus = [a for a in arquivos if a.stem in ids]
    assert len(meus) == 2, f"esperava 2 arquivos distintos, achei {arquivos}"
    assert not any(a.name == "None.pptx" for a in arquivos)

    listagem = c.get("/api/pptx-templates",
                     headers={"Authorization": f"Bearer {token}"}).json()
    assert len(listagem) == 2
    assert all(item["available"] for item in listagem)


def test_deleting_one_template_keeps_the_other_file():
    c = SyncClient()
    token = _token(c, "B")
    manter = _upload(c, token, "modelo_c").json()
    apagado = _upload(c, token, "modelo_d").json()

    resposta = c.delete(f"/api/pptx-templates/{apagado['id']}",
                        headers={"Authorization": f"Bearer {token}"})
    assert resposta.status_code == 204

    arquivos = {a.stem for a in _templates_dir().rglob("*.pptx")}
    assert apagado["id"] not in arquivos          # o alvo saiu
    assert manter["id"] in arquivos               # o outro sobreviveu

    restantes = c.get("/api/pptx-templates",
                      headers={"Authorization": f"Bearer {token}"}).json()
    assert [item["id"] for item in restantes] == [manter["id"]]
    assert restantes[0]["available"] is True


def test_listing_flags_a_template_whose_file_vanished():
    """Registro órfão precisa ser visível: a geração por mutação depende do
    arquivo original."""
    c = SyncClient()
    token = _token(c, "C")
    criado = _upload(c, token, "modelo_e").json()

    caminho = next(_templates_dir().rglob(f"{criado['id']}.pptx"))
    caminho.unlink()

    listagem = c.get("/api/pptx-templates",
                     headers={"Authorization": f"Bearer {token}"}).json()
    alvo = next(item for item in listagem if item["id"] == criado["id"])
    assert alvo["available"] is False


# ── Marcação de slots (Etapa 2) ──────────────────────────────────────────────

def _layout(c, token, template_id):
    return c.get(f"/api/pptx-templates/{template_id}",
                 headers={"Authorization": f"Bearer {token}"}).json()


def test_get_template_returns_the_layout_with_suggested_slots():
    c = SyncClient()
    token = _token(c, "D")
    criado = _upload(c, token, "modelo_f").json()

    detalhe = _layout(c, token, criado["id"])
    elementos = detalhe["layout"]["slides"][0]["elements"]
    assert all(e.get("slot") for e in elementos)
    assert all(e.get("src_shape_id") for e in elementos)


def test_patch_slots_only_changes_the_slot():
    """Geometria e formatação vêm do .pptx original — a tela de marcação não
    pode alterá-las."""
    c = SyncClient()
    token = _token(c, "E")
    criado = _upload(c, token, "modelo_g").json()
    antes = _layout(c, token, criado["id"])["layout"]["slides"][0]
    alvo = antes["elements"][0]

    resposta = c.patch(
        f"/api/pptx-templates/{criado['id']}/slots",
        headers={"Authorization": f"Bearer {token}"},
        json={"slots": {antes["id"]: {alvo["id"]: "body"}}},
    )
    assert resposta.status_code == 200, resposta.text

    depois = resposta.json()["layout"]["slides"][0]["elements"][0]
    assert depois["slot"] == "body"
    for campo in ("x", "y", "w", "h", "font_size", "color", "text"):
        assert depois.get(campo) == alvo.get(campo), f"{campo} não podia mudar"

    # persistiu?
    relido = _layout(c, token, criado["id"])["layout"]["slides"][0]["elements"][0]
    assert relido["slot"] == "body"


def test_patch_slots_rejects_unknown_slot():
    c = SyncClient()
    token = _token(c, "F")
    criado = _upload(c, token, "modelo_h").json()
    slide = _layout(c, token, criado["id"])["layout"]["slides"][0]

    resposta = c.patch(
        f"/api/pptx-templates/{criado['id']}/slots",
        headers={"Authorization": f"Bearer {token}"},
        json={"slots": {slide["id"]: {slide["elements"][0]["id"]: "titulo_errado"}}},
    )
    assert resposta.status_code == 422
    assert "titulo_errado" in resposta.text


def test_slots_of_another_user_are_not_reachable():
    c = SyncClient()
    dono = _token(c, "G")
    criado = _upload(c, dono, "modelo_i").json()
    intruso = _token(c, "H")

    assert c.get(f"/api/pptx-templates/{criado['id']}",
                 headers={"Authorization": f"Bearer {intruso}"}).status_code == 404
    assert c.patch(f"/api/pptx-templates/{criado['id']}/slots",
                   headers={"Authorization": f"Bearer {intruso}"},
                   json={"slots": {}}).status_code == 404


def test_legacy_template_is_reimported_to_gain_slots():
    """Modelos enviados antes desta etapa não têm slot nem âncora. Sem a
    âncora o exportador por mutação não acha o shape no .pptx — então a
    releitura acontece sozinha na primeira abertura."""
    from app.models import PptxTemplate

    c = SyncClient()
    token = _token(c, "I")
    criado = _upload(c, token, "modelo_legado").json()

    # simula o layout antigo (sem slot/âncora)
    sessao = _TestSession()
    registro = sessao.query(PptxTemplate).filter(PptxTemplate.id == criado["id"]).first()
    antigo = {"slides": [{
        "id": "s0", "kind": "cover",
        "elements": [{"id": "t0", "type": "text", "x": 0.1, "y": 0.1,
                      "w": 0.5, "h": 0.1, "text": "MODELO_LEGADO", "font_size": 32}],
    }]}
    registro.layout = antigo
    sessao.commit()
    sessao.close()

    detalhe = _layout(c, token, criado["id"])
    elementos = detalhe["layout"]["slides"][0]["elements"]
    assert all(e.get("slot") for e in elementos)
    assert all(e.get("src_shape_id") for e in elementos)


def test_legacy_template_without_file_does_not_break_the_screen():
    """Sem o .pptx em disco não há como reimportar — a tela ainda abre."""
    from app.models import PptxTemplate

    c = SyncClient()
    token = _token(c, "J")
    criado = _upload(c, token, "modelo_sem_arquivo").json()
    next(_templates_dir().rglob(f"{criado['id']}.pptx")).unlink()

    sessao = _TestSession()
    registro = sessao.query(PptxTemplate).filter(PptxTemplate.id == criado["id"]).first()
    registro.layout = {"slides": [{"id": "s0", "kind": "cover", "elements": [
        {"id": "t0", "type": "text", "x": 0.1, "y": 0.1, "w": 0.5, "h": 0.1,
         "text": "SEM ARQUIVO", "font_size": 32}]}]}
    sessao.commit()
    sessao.close()

    detalhe = _layout(c, token, criado["id"])
    assert detalhe["available"] is False
    assert detalhe["layout"]["slides"][0]["elements"][0]["text"] == "SEM ARQUIVO"
