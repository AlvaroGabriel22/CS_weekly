from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.pptx_service import PptxService

BASE_DATA = {
    "title": "Weekly Report - 06/07/2026 a 12/07/2026",
    "department": "CSI",
    "period_label": "06/07–12/07",
    "summary": "Auditoria da linha 3 concluída sem NC crítica; plano do fornecedor XYZ avançou 80%.",
    "highlights": ["Linha 3 liberada sem NC crítica", "Plano XYZ 80% concluído"],
    "activities": [
        {
            "source": 1,
            "title": "Auditoria de processo",
            "date": "08/07",
            "narrative": "A verificação confirmou a estabilidade dos controles da linha piloto.",
            "impact": "Produção liberada sem restrições.",
            "facts": ["Lote impactado: L2407B", "Produção afetada: 1.850 unidades"],
            "actions": ["Reunião técnica com Engenharia", "Coleta de amostras"],
        }
    ],
    "conclusions": ["Linha 3 estável e liberada."],
    "next_steps": ["Acompanhar o plano de ação do fornecedor."],
}


def _generate(tmp_path: Path, data: dict, language: str = "pt", images: list | None = None) -> str:
    service = PptxService(output_dir=str(tmp_path))
    return service.generate(
        report_id="test-report",
        template_path=None,
        slides_config={},
        data=data,
        images=images or [],
        week_number=28,
        year=2026,
        author="Usuário Teste",
        language=language,
    )


def _all_text(presentation: Presentation) -> str:
    parts = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
    return " | ".join(parts)


def test_generates_compact_v2_deck(tmp_path: Path):
    output = _generate(tmp_path, BASE_DATA)

    assert Path(output).exists()
    presentation = Presentation(output)
    # cover + one dense content slide
    assert len(presentation.slides) == 2
    text = _all_text(presentation)
    assert "Weekly Report" in text
    assert "1. Auditoria de processo" in text
    assert "Tratativas" in text
    assert "Lote impactado: L2407B" in text
    assert "DESTAQUES" not in text
    assert "SÍNTESE" not in text


def test_renders_kpi_table_when_rows_exist(tmp_path: Path):
    data = {
        **BASE_DATA,
        "kpi_table": [
            {"kpi": "FPY", "result": "97,8% → 98,4%", "trend": "▲ Melhorou"},
            {"kpi": "FCT Failure", "result": "0,68% → 0,94%", "trend": "▼ Piorou"},
        ],
    }
    output = _generate(tmp_path, data)
    presentation = Presentation(output)
    tables = [
        shape
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_table
    ]
    assert len(tables) == 1
    assert "FPY" in _all_text(presentation)


def test_ignores_invalid_kpi_table(tmp_path: Path):
    data = {**BASE_DATA, "kpi_table": ["not-a-row", 42]}
    output = _generate(tmp_path, data)
    presentation = Presentation(output)
    tables = [
        shape
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_table
    ]
    assert not tables


def test_paginates_many_activities(tmp_path: Path):
    long_narrative = (
        "Análise detalhada da evidência mostrando desvios recorrentes no processo "
        "de montagem, com correlação direta entre turno e incidência de defeitos."
    )
    data = {
        **BASE_DATA,
        "activities": [
            {
                "source": i,
                "title": f"Atividade relevante número {i}",
                "date": "08/07",
                "narrative": long_narrative,
                "impact": "Redução mensurável de retrabalho.",
                "facts": [f"Fato crucial {i}.1", f"Fato crucial {i}.2"],
                "actions": [f"Ação corretiva {i}.1", f"Ação corretiva {i}.2"],
            }
            for i in range(1, 9)
        ],
    }
    output = _generate(tmp_path, data)
    presentation = Presentation(output)
    assert len(presentation.slides) >= 3


def test_inline_activity_images(tmp_path: Path):
    image_path = tmp_path / "evidence.png"
    Image.new("RGB", (320, 200), color=(40, 90, 180)).save(image_path)

    data = {
        **BASE_DATA,
        "activities": [
            {
                **BASE_DATA["activities"][0],
                "images": [{"path": str(image_path), "caption": "Evidência"}],
            }
        ],
    }
    output = _generate(tmp_path, data)
    presentation = Presentation(output)
    pictures = [
        shape
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    assert len(pictures) == 1
    # thumbnail, not a full-slide image
    assert pictures[0].width.inches <= 1.5


def test_unmatched_images_render_as_evidence_block(tmp_path: Path):
    image_path = tmp_path / "extra.png"
    Image.new("RGB", (300, 300), color=(200, 60, 60)).save(image_path)

    output = _generate(
        tmp_path, BASE_DATA, images=[{"path": str(image_path), "caption": ""}]
    )
    presentation = Presentation(output)
    text = _all_text(presentation)
    pictures = [
        shape
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    assert "Evidências" in text
    assert len(pictures) == 1


def test_generates_english_deck(tmp_path: Path):
    output = _generate(tmp_path, BASE_DATA, language="en")
    presentation = Presentation(output)
    text = _all_text(presentation)
    assert "HIGHLIGHTS" not in text
    assert "NEXT STEPS" not in text
    assert "Actions taken" in text
