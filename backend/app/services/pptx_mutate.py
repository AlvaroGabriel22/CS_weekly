"""Primitivas de MUTAÇÃO de um .pptx existente.

Este é o coração da geração fiel ao template: em vez de reconstruir uma
apresentação do zero (perdendo tudo que o modelo interno não descreve — tema,
master, gradiente, sombra, SmartArt), abrimos o arquivo do usuário e apenas
trocamos o conteúdo dos slots. O que não é tocado permanece idêntico.

ATENÇÃO — este é o ÚNICO módulo autorizado a usar API privada do python-pptx
(`rels._rels`, `_add_relationship`, `_element`, `_spTree`). A série 1.x muda
essas internas; mantendo o uso concentrado aqui, um upgrade futuro toca um
arquivo só. Se precisar de uma mutação nova, adicione-a AQUI e exponha uma
função pública, em vez de mexer em XML no chamador.

Duas armadilhas verificadas em teste (não remova as proteções):

1. Os rIds NÃO podem ser reaproveitados. `_add_relationship` gera o rId
   sozinho, e o XML copiado ainda aponta para os rIds da origem. Por isso
   `duplicate_slide` remapeia `r:id`/`r:embed`/`r:link` na árvore copiada —
   sem isso, imagem e gráfico do slide novo apontam para o lugar errado.
2. O slide duplicado COMPARTILHA a part do gráfico com o original. Escrever
   dados em um alteraria os dois. Por isso `set_chart_data` clona a ChartPart
   antes de escrever.
"""
from __future__ import annotations

# python-pptx 0.6.x acessa collections.abc por um caminho que quebra no
# Python 3.12 se o módulo ainda não tiver sido importado. Garante isso.
import collections.abc  # noqa: F401
import copy
import io
import logging
from pathlib import Path
from typing import Any, Sequence

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.opc.packuri import PackURI
from pptx.oxml.ns import qn
from pptx.parts.chart import ChartPart
from pptx.util import Emu

logger = logging.getLogger(__name__)

R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
C_NS = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
CHART_RELTYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart"
)
# Referências a relacionamento que aparecem dentro do XML de um shape.
_REL_ATTRS = ("id", "embed", "link")


class PptxMutateError(Exception):
    """Falha de mutação com mensagem pronta para o usuário (PT)."""


# ───────────────────────────── arquivo ───────────────────────────────────────

def open_presentation(path: str | Path) -> Presentation:
    """Abre o .pptx do usuário. Erro claro se o arquivo não abrir."""
    try:
        return Presentation(str(path))
    except Exception as error:
        raise PptxMutateError(
            "Não foi possível abrir o modelo (.pptx). Envie o arquivo novamente."
        ) from error


def save(presentation: Presentation, path: str | Path) -> str:
    destination = Path(path)
    destination.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(destination))
    return str(destination)


# ───────────────────────────── slides ────────────────────────────────────────

def duplicate_slide(presentation: Presentation, index: int):
    """Duplica o slide `index` (vai para o fim) e devolve o slide novo.

    Copia shapes, relacionamentos e o fundo próprio do slide, se houver.
    """
    slides = list(presentation.slides)
    if not 0 <= index < len(slides):
        raise PptxMutateError("Slide de modelo inexistente.")
    source = slides[index]
    dest = presentation.slides.add_slide(source.slide_layout)

    # add_slide traz os placeholders do layout — o clone não os quer.
    for shape in list(dest.shapes):
        shape._element.getparent().remove(shape._element)

    # 1) relacionamentos primeiro: precisamos do mapa antigo → novo.
    rid_map: dict[str, str] = {}
    for rid, rel in source.part.rels._rels.items():
        if "slideLayout" in rel.reltype:
            continue  # o layout já veio do add_slide
        rid_map[rid] = dest.part.rels._add_relationship(
            rel.reltype, rel._target, rel.is_external
        )

    # 2) fundo próprio do slide (se o template definir um)
    background = source._element.find(qn("p:cSld") + "/" + qn("p:bg"))
    if background is None:
        c_sld = source._element.find(qn("p:cSld"))
        background = c_sld.find(qn("p:bg")) if c_sld is not None else None
    if background is not None:
        dest_csld = dest._element.find(qn("p:cSld"))
        dest_csld.insert(0, _remap_rels(copy.deepcopy(background), rid_map))

    # 3) shapes, com as referências reescritas
    for shape in source.shapes:
        element = _remap_rels(copy.deepcopy(shape._element), rid_map)
        dest.shapes._spTree.insert_element_before(element, "p:extLst")

    return dest


def _remap_rels(element, rid_map: dict[str, str]):
    """Reescreve r:id/r:embed/r:link do XML copiado para os rIds do destino."""
    if not rid_map:
        return element
    for node in element.iter():
        for attr in _REL_ATTRS:
            key = R_NS + attr
            current = node.attrib.get(key)
            if current is not None and current in rid_map:
                node.attrib[key] = rid_map[current]
    return element


def delete_slide(presentation: Presentation, index: int) -> None:
    """Remove o slide e o relacionamento correspondente."""
    id_list = presentation.slides._sldIdLst
    entries = list(id_list)
    if not 0 <= index < len(entries):
        raise PptxMutateError("Slide inexistente para remoção.")
    entry = entries[index]
    rid = entry.get(R_NS + "id")
    if rid:
        presentation.part.drop_rel(rid)
    id_list.remove(entry)


def move_slide(presentation: Presentation, old_index: int, new_index: int) -> None:
    """Reposiciona um slide na apresentação."""
    id_list = presentation.slides._sldIdLst
    entries = list(id_list)
    if not 0 <= old_index < len(entries) or not 0 <= new_index < len(entries):
        raise PptxMutateError("Posição de slide inválida.")
    entry = entries[old_index]
    id_list.remove(entry)
    id_list.insert(new_index, entry)


def slide_count(presentation: Presentation) -> int:
    return len(presentation.slides._sldIdLst)


# ───────────────────────────── texto ─────────────────────────────────────────

def set_text(shape, text: str) -> None:
    """Troca o texto MANTENDO a formatação (fonte, corpo, cor, negrito).

    A formatação de um PPT real mora quase sempre no run/parágrafo/master, não
    num atributo que saibamos reproduzir — então preservamos o primeiro run e
    apenas trocamos o conteúdo dele.
    """
    frame = _text_frame(shape)
    paragraphs = frame.paragraphs
    _set_paragraph_text(paragraphs[0], text)
    for paragraph in paragraphs[1:]:
        paragraph._p.getparent().remove(paragraph._p)


def set_paragraphs(shape, lines: Sequence[str], bullet_prefix: str = "") -> None:
    """Escreve várias linhas herdando o estilo do primeiro parágrafo."""
    cleaned = [str(line) for line in lines if str(line).strip()]
    if not cleaned:
        clear_text(shape)
        return
    frame = _text_frame(shape)
    template_p = copy.deepcopy(frame.paragraphs[0]._p)

    _set_paragraph_text(frame.paragraphs[0], f"{bullet_prefix}{cleaned[0]}")
    for paragraph in frame.paragraphs[1:]:
        paragraph._p.getparent().remove(paragraph._p)

    body = frame._txBody
    for line in cleaned[1:]:
        new_p = copy.deepcopy(template_p)
        body.append(new_p)
        _set_paragraph_text(frame.paragraphs[-1], f"{bullet_prefix}{line}")


def clear_text(shape) -> None:
    """Esvazia a caixa preservando-a (usado em slot sem conteúdo na semana)."""
    set_text(shape, "")


def _text_frame(shape):
    if not getattr(shape, "has_text_frame", False):
        raise PptxMutateError("O elemento selecionado não aceita texto.")
    return shape.text_frame


def _set_paragraph_text(paragraph, text: str) -> None:
    """Texto do parágrafo preservando a formatação do primeiro run."""
    runs = paragraph.runs
    if runs:
        runs[0].text = text
        for run in runs[1:]:
            run._r.getparent().remove(run._r)
    else:
        # Parágrafo sem run: o novo run herda o estilo do parágrafo/placeholder.
        paragraph.add_run().text = text


# ───────────────────────────── shapes ────────────────────────────────────────

def remove_shape(shape) -> None:
    """Remove o shape do slide (slot que a semana não preencheu)."""
    shape._element.getparent().remove(shape._element)


# ───────────────────────────── imagem ────────────────────────────────────────

def swap_image(picture, source: str | Path | bytes, contain: bool = True) -> None:
    """Troca a imagem preservando a moldura do template.

    `contain=True` reajusta largura/altura dentro da moldura original para
    manter a proporção da imagem nova, centralizada — evita a foto esticada.
    """
    blob = Path(source).read_bytes() if isinstance(source, (str, Path)) else source
    try:
        image_part, rid = picture.part.get_or_add_image_part(io.BytesIO(blob))
    except Exception as error:
        raise PptxMutateError("Não foi possível usar esta imagem.") from error
    picture._element.blipFill.blip.rEmbed = rid

    if not contain:
        return
    try:
        img_w, img_h = image_part.image.size
    except Exception:
        return
    if not img_w or not img_h:
        return
    box_l, box_t = picture.left, picture.top
    box_w, box_h = picture.width, picture.height
    if not box_w or not box_h:
        return
    if img_w / img_h > box_w / box_h:
        draw_w, draw_h = box_w, int(box_w * img_h / img_w)
    else:
        draw_h, draw_w = box_h, int(box_h * img_w / img_h)
    picture.width, picture.height = Emu(draw_w), Emu(draw_h)
    picture.left = Emu(int(box_l + (box_w - draw_w) / 2))
    picture.top = Emu(int(box_t + (box_h - draw_h) / 2))


# ───────────────────────────── tabela ────────────────────────────────────────

def fill_table(
    graphic_frame,
    columns: Sequence[str],
    rows: Sequence[Sequence[Any]],
) -> dict[str, int]:
    """Preenche a tabela do template preservando o estilo dela.

    As linhas de dados são clonadas a partir da ÚLTIMA linha existente (que
    carrega a formatação de corpo do template) e as sobrando são removidas.
    Devolve um relatório do que não coube, para o chamador avisar em vez de
    perder conteúdo em silêncio.
    """
    if not getattr(graphic_frame, "has_table", False):
        raise PptxMutateError("O elemento selecionado não é uma tabela.")
    table = graphic_frame.table
    n_cols = len(table.columns)
    if n_cols == 0:
        raise PptxMutateError("A tabela do modelo não tem colunas.")

    dropped_cols = max(0, len(columns) - n_cols)
    _fill_row(table.rows[0], columns[:n_cols])

    _resize_body(table, len(rows))
    for index, row in enumerate(rows, start=1):
        if index >= len(table.rows):
            break
        _fill_row(table.rows[index], [str(c) for c in row][:n_cols])

    dropped_rows = max(0, len(rows) - (len(table.rows) - 1))
    return {"dropped_rows": dropped_rows, "dropped_columns": dropped_cols}


def _fill_row(row, values: Sequence[Any]) -> None:
    for index, cell in enumerate(row.cells):
        text = str(values[index]) if index < len(values) else ""
        _set_paragraph_text(cell.text_frame.paragraphs[0], text)
        for paragraph in cell.text_frame.paragraphs[1:]:
            paragraph._p.getparent().remove(paragraph._p)


def _resize_body(table, wanted: int) -> None:
    """Ajusta a quantidade de linhas de dados clonando a última linha."""
    tbl = table._tbl
    current = len(table.rows) - 1  # sem o cabeçalho
    if current < 1:
        return  # tabela só com cabeçalho: sem linha-modelo para clonar
    if wanted > current:
        # A coleção de linhas do python-pptx não aceita índice negativo.
        model = copy.deepcopy(table.rows[len(table.rows) - 1]._tr)
        for _ in range(wanted - current):
            tbl.append(copy.deepcopy(model))
    elif wanted < current:
        for row in list(table.rows)[wanted + 1:]:
            tbl.remove(row._tr)


# ───────────────────────────── gráfico ───────────────────────────────────────

def set_chart_data(
    graphic_frame,
    slide,
    categories: Sequence[str],
    series: Sequence[tuple[str, Sequence[float]]],
) -> None:
    """Troca os dados do gráfico NATIVO preservando tipo, cores e estilo.

    Clona a part do gráfico antes de escrever: um slide duplicado compartilha
    a part com o original, e escrever direto alteraria os dois.
    """
    if not getattr(graphic_frame, "has_chart", False):
        raise PptxMutateError("O elemento selecionado não é um gráfico.")
    if not categories or not series:
        raise PptxMutateError("Sem dados para preencher o gráfico.")

    _clone_chart_part(graphic_frame, slide)

    data = CategoryChartData()
    data.categories = [str(c) for c in categories]
    for name, values in series:
        data.add_series(str(name), tuple(float(v) for v in values))
    graphic_frame.chart.replace_data(data)


def _clone_chart_part(graphic_frame, slide) -> None:
    """Dá ao gráfico deste slide uma part exclusiva."""
    chart_ref = graphic_frame._element.find(f".//{C_NS}chart")
    if chart_ref is None:
        raise PptxMutateError("Gráfico do modelo ilegível.")
    old_rid = chart_ref.get(R_NS + "id")
    slide_part = slide.part
    old_part = slide_part.rels[old_rid].target_part
    package = slide_part.package

    taken = {str(part.partname) for part in package.iter_parts()}
    index = 1
    while f"/ppt/charts/chart{index}.xml" in taken:
        index += 1

    new_part = ChartPart.load(
        PackURI(f"/ppt/charts/chart{index}.xml"),
        old_part.content_type,
        package,
        old_part.blob,
    )
    # A part do gráfico tem os próprios relacionamentos (planilha embutida,
    # cores, estilo). Sem eles o PowerPoint reclama do arquivo.
    for _, rel in old_part.rels._rels.items():
        new_part.rels._add_relationship(rel.reltype, rel._target, rel.is_external)

    chart_ref.set(R_NS + "id", slide_part.rels._add_relationship(CHART_RELTYPE, new_part))
