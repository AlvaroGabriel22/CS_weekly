"""Executive PPTX builder following the "template v2" visual language.

Design principles (extracted from frontend/template v2.pptx):
- Arial, small font sizes (8-12pt), no visible boxes/cards.
- Dense two-column content slides: activities on the left, a sidebar on the
  right (synthesis, KPI table, highlights, conclusions, next steps).
- Small blue accent bars mark section headings.
- Images are small inline thumbnails (visual reference only).
- A dynamic layout engine measures text and paginates automatically.
"""

import logging
import math
import re
from collections import deque
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.core.config import get_settings
from app.services.pptx.charts import render_chart
from app.services.pptx.profiles import resolve_layout_profile
from app.services.pptx.strings import (
    BRAND,
    CHAR_W_FACTOR,
    CONTENT_BOTTOM,
    CONTENT_TOP,
    DIVIDER_X,
    FONT,
    FOOTER_Y,
    FULL_W,
    GREEN,
    INK,
    LEFT_W,
    LINE,
    LINE_H_FACTOR,
    MARGIN_L,
    MAX_CONTENT_SLIDES,
    MUTED,
    PAGE_H,
    PAGE_W,
    RED,
    RIGHT_W,
    RIGHT_X,
    STRINGS,
    WHITE,
)


logger = logging.getLogger(__name__)
settings = get_settings()

PLACEHOLDER_RE = re.compile(r"\{\{(\w+)\}\}")


class PptxService:
    """Builds an executive QWI presentation from structured weekly content."""

    def __init__(self, output_dir: str | None = None):
        upload_root = Path(settings.UPLOAD_DIR)
        self.output_dir = Path(output_dir) if output_dir else upload_root / "reports"
        self.templates_dir = upload_root / "templates"
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.templates_dir.mkdir(parents=True, exist_ok=True)

    # ------------------------------------------------------------- legacy API

    def analyze_template(self, file_path: str) -> dict:
        """Legacy/internal template support retained for future phases."""
        presentation = Presentation(file_path)
        slides = []
        all_placeholders: set[str] = set()
        for index, slide in enumerate(presentation.slides):
            placeholders: list[str] = []
            title = ""
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text
                found = PLACEHOLDER_RE.findall(text)
                placeholders.extend(found)
                all_placeholders.update(found)
                if not title and text.strip():
                    title = text.strip()[:80]
            slides.append(
                {
                    "number": index + 1,
                    "title": title or f"Slide {index + 1}",
                    "placeholders": list(dict.fromkeys(placeholders)),
                }
            )
        return {
            "slides": slides,
            "placeholders": sorted(all_placeholders),
            "slide_count": len(slides),
        }

    def save_template_file(
        self, template_id: str, content: bytes, original_name: str
    ) -> str:
        extension = Path(original_name).suffix.lower() or ".pptx"
        path = self.templates_dir / f"{template_id}{extension}"
        path.write_bytes(content)
        return str(path)

    # ------------------------------------------------------------- generation

    def generate(
        self,
        report_id: str,
        template_path: str | None,
        slides_config: dict,
        data: dict[str, Any],
        images: list[dict],
        week_number: int,
        year: int,
        author: str,
        language: str = "pt",
    ) -> str:
        """Generate the executive deck. `images` are evidence thumbnails that
        could not be matched to a specific activity."""
        del template_path, slides_config
        strings = STRINGS.get(language, STRINGS["pt"])
        self._strings = strings
        output_path = self.output_dir / f"weekly_{report_id}.pptx"

        presentation = Presentation()
        presentation.slide_width = Inches(PAGE_W)
        presentation.slide_height = Inches(PAGE_H)

        self._add_cover(presentation, data, week_number, year, author, strings)
        layout_profile = resolve_layout_profile(
            data.get("presentation_plan"),
            data.get("sector", "CSI"),
            data.get("activities") if isinstance(data.get("activities"), list) else [],
        )
        self._layout_profile = layout_profile
        self._render_global_blocks(presentation, data, strings, week_number, year)
        self._render_content_slides(
            presentation, data, images or [], strings, week_number, year, layout_profile
        )
        presentation.save(str(output_path))

        logger.info("PPTX generated | path=%s", output_path)
        return str(output_path)

    def _blank_slide(self, presentation: Presentation):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = WHITE
        return slide

    # ------------------------------------------------------------------ cover

    def _add_cover(
        self,
        presentation: Presentation,
        data: dict[str, Any],
        week_number: int,
        year: int,
        author: str,
        strings: dict[str, str],
    ) -> None:
        slide = self._blank_slide(presentation)
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), presentation.slide_height
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = BRAND
        accent.line.fill.background()

        label = self._text_box(
            slide, 0.9, 0.8, 3.5, 0.4, "QUALITY WEEKLY INTELLIGENCE", 11, BRAND, bold=True
        )
        del label

        title = str(data.get("title") or "Weekly Report")
        self._text_box(slide, 0.9, 2.0, 10.8, 1.3, title, 32, INK, bold=True)

        department = str(data.get("department") or "")
        subtitle = f"{strings['week']} {week_number} · {year}"
        if department:
            subtitle += f"  |  {department}"
        self._text_box(slide, 0.95, 3.45, 10.5, 0.5, subtitle, 16, MUTED)

        self._text_box(slide, 0.95, 5.95, 9.0, 0.4, author, 13, INK)
        self._text_box(
            slide, 0.95, 6.35, 9.0, 0.35, strings["cover_tagline"], 10, MUTED
        )
        self._page_number(slide, 1)

    # -------------------------------------------------------- content layout

    def _render_global_blocks(
        self,
        presentation: Presentation,
        data: dict[str, Any],
        strings: dict[str, str],
        week_number: int,
        year: int,
    ) -> None:
        plan = data.get("presentation_plan") or {}
        blocks = plan.get("global_blocks") or []
        if not blocks:
            return
        week_label = f"W{week_number}"
        period_label = self._as_text(data.get("period_label")) or (
            f"{strings['week']} {week_number} · {year}"
        )
        slide = self._blank_slide(presentation)
        self._content_header(slide, week_label, period_label)
        y = CONTENT_TOP
        for block in blocks:
            if not isinstance(block, dict):
                continue
            if block.get("type") == "chart":
                height = 2.8
                y += render_chart(slide, block, MARGIN_L, y, FULL_W + 2.2, height)
            elif block.get("type") in ("generic_table", "measurement_table"):
                y += self._render_data_table(slide, block, MARGIN_L, y, FULL_W + 2.2, strings)
            y += 0.12
        self._footer(slide, len(presentation.slides), strings)

    def _render_content_slides(
        self,
        presentation: Presentation,
        data: dict[str, Any],
        unmatched_images: list[dict],
        strings: dict[str, str],
        week_number: int,
        year: int,
        layout_profile: str = "executive",
    ) -> None:
        left_queue: deque[list[dict]] = deque()
        activities = data.get("activities")
        if not isinstance(activities, list):
            activities = []

        kpi_rows = self._valid_kpi_rows(data.get("kpi_table"))
        if kpi_rows:
            left_queue.append(
                [
                    {"t": "caption", "text": strings["kpi_table"]},
                    {"t": "table", "rows": kpi_rows},
                ]
            )

        for index, item in enumerate(activities, start=1):
            left_queue.append(self._activity_block(index, item, strings))

        if unmatched_images:
            left_queue.append(self._evidence_block(unmatched_images, strings))

        week_label = f"W{week_number}"
        period_label = self._as_text(data.get("period_label")) or (
            f"{strings['week']} {week_number} · {year}"
        )

        slides_rendered = 0
        while left_queue and slides_rendered < MAX_CONTENT_SLIDES:
            slide = self._blank_slide(presentation)
            slides_rendered += 1
            self._content_header(slide, week_label, period_label)

            left_width = FULL_W
            y = CONTENT_TOP
            while left_queue:
                block = left_queue[0]
                fitted, remainder = self._split_block(
                    block, left_width, CONTENT_BOTTOM - y
                )
                if not fitted:
                    break
                y = self._render_block(slide, fitted, MARGIN_L, y, left_width)
                if remainder:
                    left_queue[0] = remainder
                    break
                left_queue.popleft()

            self._footer(slide, len(presentation.slides), strings)

    def _build_right_queue(
        self, data: dict[str, Any], strings: dict[str, str], sidebar_keys: list[str]
    ) -> deque:
        queue: deque[list[dict]] = deque()

        if "synthesis" in sidebar_keys:
            summary = self._as_text(data.get("summary"))
            if summary:
                queue.append(
                    [
                        {"t": "caption", "text": strings["synthesis"]},
                        {"t": "text", "text": summary, "size": 10, "color": INK},
                    ]
                )

        if "kpi_table" in sidebar_keys:
            kpi_rows = self._valid_kpi_rows(data.get("kpi_table"))
            if kpi_rows:
                queue.append(
                    [
                        {"t": "caption", "text": strings["kpi_table"]},
                        {"t": "table", "rows": kpi_rows},
                    ]
                )

        if "highlights" in sidebar_keys:
            highlights = self._as_list(data.get("highlights"))[:4]
            if highlights:
                first, *rest = highlights
                queue.append(
                    [
                        {"t": "caption", "text": strings["highlights"]},
                        {"t": "hl", "text": first},
                    ]
                )
                for item in rest:
                    queue.append([{"t": "hl", "text": item}])

        if "conclusions" in sidebar_keys:
            conclusions = self._as_list(data.get("conclusions"))[:3]
            if conclusions:
                first, *rest = conclusions
                queue.append(
                    [
                        {"t": "caption", "text": strings["conclusions"]},
                        {"t": "text", "text": f"•  {first}", "size": 10, "color": INK},
                    ]
                )
                for item in rest:
                    queue.append(
                        [{"t": "text", "text": f"•  {item}", "size": 10, "color": INK}]
                    )

        if "next_steps" in sidebar_keys:
            next_steps = self._as_list(data.get("next_steps"))[:3]
            if next_steps:
                first, *rest = next_steps
                queue.append(
                    [
                        {"t": "caption", "text": strings["next_steps"]},
                        {"t": "text", "text": f"•  {first}", "size": 10, "color": INK},
                    ]
                )
                for item in rest:
                    queue.append(
                        [{"t": "text", "text": f"•  {item}", "size": 10, "color": INK}]
                    )

        return queue

    def _activity_block(
        self, index: int, item: Any, strings: dict[str, str]
    ) -> list[dict]:
        if not isinstance(item, dict):
            item = {"title": str(item)}

        elements: list[dict] = []
        title = self._as_text(item.get("title")) or "—"
        elements.append(
            {
                "t": "heading",
                "text": f"{index}. {title}",
                "date": self._as_text(item.get("date")),
            }
        )

        narrative = self._as_text(item.get("narrative"))
        content_mode = self._as_text(item.get("content_mode")) or "compress"
        if narrative:
            max_chars = 400 if content_mode == "transcribe" else 220
            if len(narrative) > max_chars:
                narrative = narrative[: max_chars - 1].rsplit(" ", 1)[0] + "…"
            elements.append(
                {"t": "text", "text": narrative, "size": 10, "color": MUTED}
            )

        for fact in self._as_list(item.get("facts"))[:4]:
            elements.append({"t": "text", "text": fact, "size": 10, "color": INK})

        impact = self._as_text(item.get("impact"))
        if impact and not self._is_empty_impact(impact):
            elements.append(
                {
                    "t": "text",
                    "text": f"{strings['impact']}: {impact}",
                    "size": 10,
                    "color": INK,
                    "bold": True,
                }
            )

        actions = self._as_list(item.get("actions"))[:4]
        if actions:
            elements.append({"t": "label", "text": strings["treatments"]})
            for action in actions:
                elements.append(
                    {"t": "text", "text": f"•  {action}", "size": 10, "color": INK}
                )

        images = [
            img
            for img in (item.get("images") or [])
            if isinstance(img, dict) and Path(str(img.get("path", ""))).exists()
        ]
        if images:
            elements.append({"t": "images", "images": images[:3]})

        for block in item.get("blocks") or []:
            if isinstance(block, dict):
                elements.extend(self._visual_block_elements(block, strings))

        elements.append({"t": "gap", "h": 0.26})
        return elements

    def _visual_block_elements(
        self, block: dict[str, Any], strings: dict[str, str]
    ) -> list[dict]:
        block_type = block.get("type")
        if block_type == "device_info":
            return [{"t": "device_info", "block": block, "label": strings.get("device_info", "Identificação")}]
        if block_type in ("measurement_table", "generic_table"):
            title = block.get("title") or strings.get("measurements", "Medições")
            return [{"t": "data_table", "block": block, "label": title}]
        if block_type == "countermeasure_table":
            return [
                {
                    "t": "countermeasure",
                    "block": block,
                    "label": strings.get("countermeasures", "Contramedidas"),
                }
            ]
        if block_type == "chart":
            return [{"t": "chart", "block": block}]
        if block_type == "image_row":
            imgs = []
            for ref in block.get("images") or []:
                if isinstance(ref, dict) and ref.get("path") and Path(str(ref["path"])).exists():
                    imgs.append(ref)
            if imgs:
                return [{"t": "images", "images": imgs[:3]}]
        if block_type in ("text", "highlight") and block.get("text"):
            return [
                {
                    "t": "text",
                    "text": block["text"],
                    "size": 10,
                    "color": INK,
                    "bold": block_type == "highlight",
                }
            ]
        return []

    def _is_empty_impact(self, impact: str) -> bool:
        """The model sometimes writes the 'no impact' instruction literally."""
        normalized = impact.lower().strip(" .")
        return normalized in {
            "no measured impact",
            "sem impacto medido",
            "sem impacto mensurado",
            "nenhum impacto medido",
            "nenhum impacto mensurado",
            "n/a",
            "none",
        }

    def _evidence_block(
        self, images: list[dict], strings: dict[str, str]
    ) -> list[dict]:
        existing = [
            img for img in images if Path(str(img.get("path", ""))).exists()
        ]
        elements: list[dict] = [{"t": "label", "text": strings["evidences"]}]
        for start in range(0, min(len(existing), 6), 3):
            elements.append({"t": "images", "images": existing[start : start + 3]})
        elements.append({"t": "gap", "h": 0.2})
        return elements

    # ----------------------------------------------------- element rendering

    def _measure_block(self, block: list[dict], width: float) -> float:
        return sum(self._measure_element(el, width) for el in block)

    def _split_block(
        self, block: list[dict], width: float, available: float
    ) -> tuple[list[dict], list[dict]]:
        """Fit as many elements as possible into `available` inches.

        Returns (fitted, remainder). Headings/labels are never left orphaned at
        the bottom of a column; continuation slides repeat the heading with a
        '(cont.)' suffix so the reader keeps context.
        """
        fitted: list[dict] = []
        used = 0.0
        for index, element in enumerate(block):
            height = self._measure_element(element, width)
            if used + height > available:
                remainder = block[index:]
                # move an orphaned heading/label down with its content
                while fitted and fitted[-1]["t"] in ("heading", "label"):
                    remainder.insert(0, fitted.pop())
                if not fitted and available >= CONTENT_BOTTOM - CONTENT_TOP - 0.01:
                    # element taller than a full column: render it anyway
                    fitted = block[: index + 1]
                    remainder = block[index + 1 :]
                if fitted and remainder:
                    remainder = self._with_continuation_heading(block, remainder)
                if all(el["t"] == "gap" for el in remainder):
                    remainder = []
                return fitted, remainder
            fitted.append(element)
            used += height
        return fitted, []

    def _with_continuation_heading(
        self, block: list[dict], remainder: list[dict]
    ) -> list[dict]:
        if remainder and remainder[0]["t"] in ("heading", "label"):
            return remainder
        heading = next((el for el in block if el["t"] == "heading"), None)
        if heading is None:
            return remainder
        continuation = {
            "t": "heading",
            "text": f"{heading['text']} (cont.)",
            "date": None,
        }
        return [continuation, *remainder]

    def _measure_element(self, el: dict, width: float) -> float:
        kind = el["t"]
        if kind == "heading":
            return max(0.32, self._text_height(el["text"], width - 1.15, 12) + 0.05)
        if kind == "label":
            return 0.32
        if kind == "caption":
            return 0.28
        if kind == "text":
            return self._text_height(el["text"], width, el.get("size", 10)) + 0.02
        if kind == "hl":
            return self._text_height(el["text"], width - 0.22, 10.5) + 0.1
        if kind == "images":
            return 1.02
        if kind == "table":
            return 0.34 + 0.27 * len(el["rows"])
        if kind == "device_info":
            fields = (el.get("block") or {}).get("fields") or {}
            return 0.34 + 0.22 * max(1, math.ceil(len(fields) / 2))
        if kind == "data_table":
            rows = (el.get("block") or {}).get("rows") or []
            return 0.38 + 0.26 * min(len(rows), 8)
        if kind == "countermeasure":
            rows = (el.get("block") or {}).get("rows") or []
            return 0.38 + 0.26 * min(len(rows), 6)
        if kind == "chart":
            return 2.9
        if kind == "gap":
            return el["h"]
        return 0.0

    def _render_block(
        self, slide, block: list[dict], x: float, y: float, width: float
    ) -> float:
        for el in block:
            y += self._render_element(slide, el, x, y, width)
        return y

    def _render_element(
        self, slide, el: dict, x: float, y: float, width: float
    ) -> float:
        kind = el["t"]
        height = self._measure_element(el, width)

        if kind == "heading":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x - 0.34),
                Inches(y + 0.075),
                Inches(0.28),
                Inches(0.09),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = BRAND
            bar.line.fill.background()
            self._text_box(
                slide, x, y, width - 1.15, height, el["text"], 12, INK, bold=True
            )
            date = el.get("date")
            if date:
                box = self._text_box(
                    slide, x + width - 1.0, y + 0.03, 1.0, 0.22, date, 9, MUTED
                )
                box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        elif kind == "label":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x - 0.34),
                Inches(y + 0.07),
                Inches(0.28),
                Inches(0.09),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = BRAND
            bar.line.fill.background()
            self._text_box(slide, x, y, width, 0.28, el["text"], 11, INK, bold=True)

        elif kind == "caption":
            self._text_box(
                slide, x, y, width, 0.22, str(el["text"]).upper(), 9, BRAND, bold=True
            )

        elif kind == "text":
            self._multiline_text(
                slide,
                x,
                y,
                width,
                height,
                el["text"],
                el.get("size", 10),
                el.get("color", INK),
                bold=el.get("bold", False),
            )

        elif kind == "hl":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x),
                Inches(y + 0.02),
                Inches(0.045),
                Inches(max(0.2, height - 0.12)),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = BRAND
            bar.line.fill.background()
            self._multiline_text(
                slide, x + 0.18, y, width - 0.22, height, el["text"], 10.5, INK, bold=True
            )

        elif kind == "images":
            for i, image_data in enumerate(el["images"][:3]):
                cell_x = x + i * 1.55
                try:
                    left, top, w, h = self._fit_image(
                        Path(image_data["path"]), cell_x, y + 0.08, 1.4, 0.86
                    )
                    slide.shapes.add_picture(
                        str(image_data["path"]),
                        Inches(left),
                        Inches(top),
                        width=Inches(w),
                        height=Inches(h),
                    )
                except Exception as error:
                    logger.warning(
                        "Failed to insert image %s: %s", image_data.get("path"), error
                    )

        elif kind == "table":
            self._render_kpi_table(slide, el["rows"], x, y, width)

        elif kind == "device_info":
            self._render_device_info(slide, el, x, y, width)

        elif kind == "data_table":
            self._render_data_table(slide, el.get("block") or {}, x, y, width, self._strings)

        elif kind == "countermeasure":
            self._render_countermeasure_table(
                slide, el.get("block") or {}, x, y, width, self._strings
            )

        elif kind == "chart":
            render_chart(slide, el.get("block") or {}, x, y, width, 2.6)

        return height

    def _render_device_info(
        self, slide, el: dict, x: float, y: float, width: float
    ) -> None:
        block = el.get("block") or {}
        label = el.get("label") or "Identificação"
        self._text_box(slide, x, y, width, 0.22, str(label).upper(), 9, BRAND, bold=True)
        fields = block.get("fields") or {}
        if not fields:
            return
        col_w = width / 2 - 0.1
        row_y = y + 0.28
        for index, (key, value) in enumerate(fields.items()):
            col = index % 2
            line_x = x + col * (col_w + 0.2)
            if index > 0 and index % 2 == 0:
                row_y += 0.22
            self._text_box(
                slide,
                line_x,
                row_y,
                col_w,
                0.2,
                f"{key}: {value}",
                9,
                INK,
            )

    def _render_data_table(
        self,
        slide,
        block: dict[str, Any],
        x: float,
        y: float,
        width: float,
        strings: dict[str, str],
    ) -> float:
        columns = block.get("columns") or []
        rows = block.get("rows") or []
        if not columns or not rows:
            return 0.0
        title = block.get("title")
        offset = 0.0
        if title:
            self._text_box(slide, x, y, width, 0.22, str(title).upper(), 9, BRAND, bold=True)
            offset = 0.26
        row_count = min(len(rows), 10)
        table_height = 0.28 + 0.24 * row_count
        frame = slide.shapes.add_table(
            row_count + 1,
            len(columns),
            Inches(x),
            Inches(y + offset),
            Inches(width),
            Inches(table_height),
        )
        table = frame.table
        col_w = width / max(len(columns), 1)
        for ci, header in enumerate(columns):
            table.columns[ci].width = Inches(col_w)
            cell = table.cell(0, ci)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BRAND
            para = cell.text_frame.paragraphs[0]
            para.font.name = FONT
            para.font.size = Pt(8)
            para.font.bold = True
            para.font.color.rgb = WHITE
        for ri, row in enumerate(rows[:row_count], start=1):
            for ci in range(len(columns)):
                value = row[ci] if ci < len(row) else ""
                cell = table.cell(ri, ci)
                cell.text = str(value)
                para = cell.text_frame.paragraphs[0]
                para.font.name = FONT
                para.font.size = Pt(8)
                para.font.color.rgb = INK
        return offset + table_height

    def _render_countermeasure_table(
        self,
        slide,
        block: dict[str, Any],
        x: float,
        y: float,
        width: float,
        strings: dict[str, str],
    ) -> float:
        rows = block.get("rows") or []
        if not rows:
            return 0.0
        title = block.get("title") or strings.get("countermeasures", "Contramedidas")
        self._text_box(slide, x, y, width, 0.22, str(title).upper(), 9, BRAND, bold=True)
        headers = (
            strings.get("action", "Ação"),
            strings.get("owner", "Responsável"),
            strings.get("status", "Status"),
            strings.get("due", "Prazo"),
        )
        data_rows = []
        for row in rows[:6]:
            if isinstance(row, dict):
                data_rows.append(
                    [
                        row.get("action", ""),
                        row.get("owner", ""),
                        row.get("status", ""),
                        row.get("due", ""),
                    ]
                )
        if not data_rows:
            return 0.26
        table_height = 0.28 + 0.24 * len(data_rows)
        frame = slide.shapes.add_table(
            len(data_rows) + 1,
            4,
            Inches(x),
            Inches(y + 0.26),
            Inches(width),
            Inches(table_height),
        )
        table = frame.table
        for ci, header in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = BRAND
            para = cell.text_frame.paragraphs[0]
            para.font.name = FONT
            para.font.size = Pt(8)
            para.font.bold = True
            para.font.color.rgb = WHITE
        for ri, row in enumerate(data_rows, start=1):
            for ci, value in enumerate(row):
                cell = table.cell(ri, ci)
                cell.text = str(value)
                para = cell.text_frame.paragraphs[0]
                para.font.name = FONT
                para.font.size = Pt(8)
                para.font.color.rgb = INK
        return 0.26 + table_height

    def _render_kpi_table(
        self, slide, rows: list[dict], x: float, y: float, width: float
    ) -> None:
        strings = getattr(self, "_strings", STRINGS["pt"])
        headers = (strings["kpi"], strings["result"], strings["trend"])
        table_height = 0.3 + 0.27 * len(rows)
        frame = slide.shapes.add_table(
            len(rows) + 1, 3, Inches(x), Inches(y + 0.04), Inches(width), Inches(table_height)
        )
        table = frame.table
        table.first_row = True
        try:
            table.horz_banding = False
        except Exception:
            pass

        table.columns[0].width = Inches(width * 0.42)
        table.columns[1].width = Inches(width * 0.32)
        table.columns[2].width = Inches(width * 0.26)

        for ci, header in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = BRAND
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.name = FONT
            paragraph.font.size = Pt(9)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)

        table.rows[0].height = Inches(0.3)
        for ri, row in enumerate(rows, start=1):
            table.rows[ri].height = Inches(0.27)
            trend = self._as_text(row.get("trend"))
            values = (
                self._as_text(row.get("kpi")),
                self._as_text(row.get("result")),
                trend,
            )
            for ci, value in enumerate(values):
                cell = table.cell(ri, ci)
                cell.text = value
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.name = FONT
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = INK
                cell.margin_top = Inches(0.02)
                cell.margin_bottom = Inches(0.02)
                if ci == 2:
                    if trend.startswith("▲"):
                        paragraph.font.color.rgb = GREEN
                    elif trend.startswith("▼"):
                        paragraph.font.color.rgb = RED
                    else:
                        paragraph.font.color.rgb = MUTED

    def _valid_kpi_rows(self, value: Any) -> list[dict]:
        if not isinstance(value, list):
            return []
        rows = []
        for row in value:
            if isinstance(row, dict) and self._as_text(row.get("kpi")):
                rows.append(row)
        return rows[:9]

    # -------------------------------------------------------- chrome & text

    def _content_header(self, slide, week_label: str, period_label: str) -> None:
        self._text_box(slide, MARGIN_L, 0.42, 2.5, 0.4, week_label, 18, INK, bold=True)
        box = self._text_box(
            slide, PAGE_W - 4.1, 0.52, 3.35, 0.28, period_label, 10, MUTED
        )
        box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(MARGIN_L - 0.06),
            Inches(0.88),
            Inches(6.17),
            Inches(0.075),
        )
        underline.fill.solid()
        underline.fill.fore_color.rgb = BRAND
        underline.line.fill.background()

    def _footer(self, slide, page: int, strings: dict[str, str]) -> None:
        self._text_box(
            slide, MARGIN_L, FOOTER_Y, 4.5, 0.2, strings["footer"], 8, MUTED
        )
        self._page_number(slide, page)

    def _page_number(self, slide, page: int) -> None:
        box = self._text_box(slide, 12.0, 7.02, 0.55, 0.2, str(page), 8, MUTED)
        box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    def _text_box(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        size: float,
        color: RGBColor,
        bold: bool = False,
    ):
        box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = 0
        frame.margin_right = 0
        frame.margin_top = 0
        frame.margin_bottom = 0
        paragraph = frame.paragraphs[0]
        paragraph.text = text or ""
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.font.bold = bold
        return box

    def _multiline_text(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        size: float,
        color: RGBColor,
        bold: bool = False,
    ):
        box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = 0
        frame.margin_right = 0
        frame.margin_top = 0
        frame.margin_bottom = 0
        lines = (text or "").splitlines() or [""]
        for index, line in enumerate(lines):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = line
            paragraph.font.name = FONT
            paragraph.font.size = Pt(size)
            paragraph.font.color.rgb = color
            paragraph.font.bold = bold
            paragraph.line_spacing = 1.15
        return box

    # ------------------------------------------------------------ measuring

    def _chars_per_line(self, width_in: float, size_pt: float) -> int:
        return max(12, int((width_in * 72) / (size_pt * CHAR_W_FACTOR)))

    def _text_height(self, text: str, width_in: float, size_pt: float) -> float:
        text = text or ""
        per_line = self._chars_per_line(width_in, size_pt)
        lines = 0
        for para in text.splitlines() or [""]:
            lines += max(1, math.ceil(len(para) / per_line))
        return lines * size_pt * LINE_H_FACTOR / 72 + 0.02

    # ----------------------------------------------------------- data prep

    def build_data_from_content(
        self,
        activities: list,
        ai_content: str,
        structured: dict,
        user_name: str,
        department: str,
        week_number: int,
        year: int,
        title: str | None = None,
        period_label: str | None = None,
        period_caption: str | None = None,
        sector: Any = None,
    ) -> dict[str, Any]:
        fallback_activities = []
        for index, activity in enumerate(
            sorted(activities, key=lambda item: item.activity_date), start=1
        ):
            metadata = activity.metadata_entry
            fallback_activities.append(
                {
                    "source": index,
                    "title": activity.title,
                    "date": activity.activity_date.strftime("%d/%m"),
                    "narrative": activity.description
                    or (metadata.technical_summary if metadata else "")
                    or activity.title,
                    "impact": "",
                    "facts": [],
                    "actions": [],
                }
            )

        structured_activities = structured.get("activities")
        if not isinstance(structured_activities, list) or not structured_activities:
            structured_activities = fallback_activities

        return {
            **structured,
            "title": title or f"Weekly Report · {week_number}/{year}",
            "week": f"{week_number}/{year}",
            "author": user_name,
            "department": department,
            "sector": getattr(sector, "value", sector) if sector else "CSI",
            "period_label": period_label,
            "period_caption": period_caption,
            "summary": structured.get("summary") or ai_content[:800],
            "activities": structured_activities,
            "presentation_plan": structured.get("presentation_plan")
            or {
                "layout_profile": "executive",
                "sidebar": ["synthesis", "kpi_table", "highlights", "conclusions", "next_steps"],
                "global_blocks": [],
            },
        }

    def _fit_image(
        self,
        path: Path,
        area_left: float,
        area_top: float,
        area_width: float,
        area_height: float,
    ) -> tuple[float, float, float, float]:
        with Image.open(path) as image:
            image_ratio = image.width / image.height
        area_ratio = area_width / area_height
        if image_ratio > area_ratio:
            width = area_width
            height = width / image_ratio
        else:
            height = area_height
            width = height * image_ratio
        left = area_left + (area_width - width) / 2
        top = area_top + (area_height - height) / 2
        return left, top, width, height

    def _as_list(self, value: Any) -> list[str]:
        if isinstance(value, list):
            result = []
            for item in value:
                if isinstance(item, dict):
                    result.append(
                        self._as_text(
                            item.get("title")
                            or item.get("name")
                            or item.get("summary")
                        )
                    )
                else:
                    result.append(self._as_text(item))
            return [item for item in result if item]
        if isinstance(value, str) and value.strip():
            return [
                line.strip("•- \t")
                for line in value.splitlines()
                if line.strip("•- \t")
            ]
        return []

    def _as_text(self, value: Any) -> str:
        if value is None:
            return ""
        if isinstance(value, list):
            return "\n".join(self._as_list(value))
        if isinstance(value, dict):
            return "; ".join(
                f"{key}: {self._as_text(item)}" for key, item in value.items()
            )
        return str(value).strip()
