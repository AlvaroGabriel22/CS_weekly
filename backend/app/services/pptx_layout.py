"""Renderiza o PPTX a partir do layout do editor WYSIWYG (tela de Montagem).

Contrato (espelhado em frontend/src/components/reports/slideLayout.ts):

{
  "slides": [
    {
      "id": "s1",
      "kind": "cover" | "activity" | "summary" | "kpis" | "conclusions" | "custom",
      "activity_id": "<uuid>",              # opcional (slides de atividade)
      "elements": [
        {
          "id": "e1",
          "type": "text" | "image",
          "x": 0.05, "y": 0.08, "w": 0.9, "h": 0.12,   # frações da página (0–1)
          "text": "…",                       # texto literal (editado pelo usuário)
          "binding": "summary" | "highlights" | "kpis" | "conclusions" | "next_steps",
          "attachment_id": "<uuid>",         # imagens
          "font_size": 20,                   # pt
          "bold": true,
          "align": "left" | "center" | "right",
          "color": "#0C379C",
          "pinned": false
        }
      ]
    }
  ]
}

Elementos com `binding` recebem o texto gerado pela IA no momento da geração;
elementos de texto sem binding são renderizados literalmente (liberdade total
do usuário). Posições/tamanhos são frações da página 16:9 (13.333in × 7.5in).
"""

import logging
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

PAGE_W = 13.333
PAGE_H = 7.5

_ALIGN = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


def _hex_to_rgb(value: str | None) -> RGBColor:
    try:
        v = (value or "#1F2937").lstrip("#")
        return RGBColor(int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))
    except Exception:
        return RGBColor(0x1F, 0x29, 0x37)


def _as_lines(value: Any) -> list[str]:
    """Normaliza conteúdo estruturado da IA em linhas de texto."""
    if value is None:
        return []
    if isinstance(value, str):
        return [line for line in value.splitlines() if line.strip()] or ([value] if value.strip() else [])
    if isinstance(value, list):
        lines: list[str] = []
        for item in value:
            if isinstance(item, str) and item.strip():
                lines.append(item.strip())
            elif isinstance(item, dict):
                name = item.get("name") or item.get("title") or item.get("kpi") or ""
                val = item.get("value") or item.get("target") or ""
                text = f"{name}: {val}" if name and val else (name or str(val))
                if text.strip():
                    lines.append(text.strip())
        return lines
    return [str(value)]


def resolve_binding(binding: str, structured: dict[str, Any]) -> list[str]:
    """Texto da IA para um binding. Lista de linhas (viram parágrafos/bullets)."""
    if not isinstance(structured, dict):
        return []
    key = binding.strip().lower()
    if key in {"summary", "highlights", "kpis", "conclusions", "next_steps"}:
        return _as_lines(structured.get(key))
    return []


class PptxLayoutRenderer:
    """Gera o .pptx honrando exatamente o layout do editor de montagem."""

    def __init__(self, output_dir: str | Path):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def generate(
        self,
        report_id: str,
        layout: dict[str, Any],
        structured: dict[str, Any],
        attachments: dict[str, dict[str, Any]],
    ) -> str:
        """`attachments`: {attachment_id: {"path": str, "table": dict | None}}."""
        presentation = Presentation()
        presentation.slide_width = Inches(PAGE_W)
        presentation.slide_height = Inches(PAGE_H)
        blank = presentation.slide_layouts[6]

        slides = layout.get("slides") if isinstance(layout, dict) else None
        if not isinstance(slides, list) or not slides:
            raise ValueError("Layout sem slides")

        for slide_def in slides:
            slide = presentation.slides.add_slide(blank)
            elements = slide_def.get("elements") if isinstance(slide_def, dict) else None
            for element in elements or []:
                if not isinstance(element, dict):
                    continue
                try:
                    self._render_element(slide, element, structured, attachments)
                except Exception as error:  # elemento inválido não derruba o deck
                    logger.warning(
                        "Elemento de layout ignorado | slide=%s el=%s err=%s",
                        slide_def.get("id"), element.get("id"), error,
                    )

        output_path = self.output_dir / f"weekly_{report_id}.pptx"
        presentation.save(str(output_path))
        logger.info("PPTX (layout custom) gerado | path=%s", output_path)
        return str(output_path)

    # ── elementos ──────────────────────────────────────────────────────────

    def _frame(self, element: dict[str, Any]) -> tuple[float, float, float, float]:
        # Linhas podem ter largura/altura 0 (horizontais/verticais).
        min_size = 0.0 if element.get("shape") == "line" else 0.02
        x = min(max(float(element.get("x", 0.05)), 0.0), 0.98)
        y = min(max(float(element.get("y", 0.05)), 0.0), 0.98)
        w = min(max(float(element.get("w", 0.4)), min_size), 1.0 - x)
        h = min(max(float(element.get("h", 0.1)), min_size), 1.0 - y)
        return x * PAGE_W, y * PAGE_H, w * PAGE_W, h * PAGE_H

    def _render_element(
        self,
        slide: Any,
        element: dict[str, Any],
        structured: dict[str, Any],
        attachments: dict[str, dict[str, Any]],
    ) -> None:
        kind = element.get("type")
        left, top, width, height = self._frame(element)

        if kind == "image":
            attachment_id = str(element.get("attachment_id") or "")
            path = (attachments.get(attachment_id) or {}).get("path")
            if not path or not Path(path).exists():
                logger.warning("Imagem do layout não encontrada | attachment_id=%s", attachment_id)
                return
            # A imagem NUNCA é esticada/cortada: encaixa inteira dentro da
            # caixa preservando a proporção, centralizada (regra "contain" —
            # a mesma que o editor aplica na tela).
            draw_left, draw_top, draw_w, draw_h = left, top, width, height
            try:
                from PIL import Image as PILImage

                with PILImage.open(path) as img:
                    img_w, img_h = img.size
                if img_w > 0 and img_h > 0 and width > 0 and height > 0:
                    img_ratio = img_w / img_h
                    box_ratio = width / height
                    if img_ratio > box_ratio:
                        draw_w = width
                        draw_h = width / img_ratio
                    else:
                        draw_h = height
                        draw_w = height * img_ratio
                    draw_left = left + (width - draw_w) / 2
                    draw_top = top + (height - draw_h) / 2
            except Exception as error:
                logger.warning("Falha ao medir imagem (usa caixa cheia) | %s", error)
            slide.shapes.add_picture(
                path,
                Inches(draw_left),
                Inches(draw_top),
                width=Inches(draw_w),
                height=Inches(draw_h),
            )
            return

        if kind == "table":
            attachment_id = str(element.get("attachment_id") or "")
            table_data = (attachments.get(attachment_id) or {}).get("table")
            if not table_data:
                logger.warning("Tabela do layout não encontrada | attachment_id=%s", attachment_id)
                return
            self._render_table(slide, element, table_data, left, top, width, height)
            return

        if kind == "shape":
            self._render_shape(slide, element, left, top, width, height)
            return

        # texto (literal ou binding)
        binding = element.get("binding")
        if binding:
            lines = resolve_binding(str(binding), structured)
            bullet = len(lines) > 1
        else:
            raw = str(element.get("text") or "")
            lines = raw.splitlines() or [""]
            bullet = False
        if not lines:
            return

        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.TOP
        frame.margin_left = frame.margin_right = Pt(2)
        frame.margin_top = frame.margin_bottom = Pt(1)

        size = Pt(float(element.get("font_size", 14)))
        bold = bool(element.get("bold", False))
        italic = bool(element.get("italic", False))
        font_name = str(element.get("font_family") or "Calibri")
        color = _hex_to_rgb(element.get("color"))
        align = _ALIGN.get(str(element.get("align") or "left"), PP_ALIGN.LEFT)

        for index, line in enumerate(lines):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.alignment = align
            run = paragraph.add_run()
            run.text = f"• {line}" if bullet else line
            run.font.size = size
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
            run.font.name = font_name

    def _render_table(
        self,
        slide: Any,
        element: dict[str, Any],
        table_data: dict[str, Any],
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        columns = [str(c) for c in table_data.get("columns") or []]
        rows = [[str(c) for c in r] for r in table_data.get("rows") or []]
        if not columns:
            return
        n_rows = 1 + len(rows)
        n_cols = len(columns)
        shape = slide.shapes.add_table(
            n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        table = shape.table
        font_size = Pt(float(element.get("font_size", 11)))
        header_fill = _hex_to_rgb(element.get("color") or "#0C379C")
        for c, name in enumerate(columns):
            cell = table.cell(0, c)
            cell.text = name
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = font_size
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        for r, row in enumerate(rows, start=1):
            for c in range(n_cols):
                cell = table.cell(r, c)
                cell.text = row[c] if c < len(row) else ""
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = font_size
                        run.font.color.rgb = RGBColor(0x1F, 0x29, 0x37)

    def _render_shape(
        self,
        slide: Any,
        element: dict[str, Any],
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        shape_kind = str(element.get("shape") or "rect")
        stroke = _hex_to_rgb(element.get("color") or "#0C379C")
        fill_hex = element.get("fill")

        if shape_kind == "line":
            connector = slide.shapes.add_connector(
                1,  # MSO_CONNECTOR.STRAIGHT
                Inches(left), Inches(top),
                Inches(left + width), Inches(top + height),
            )
            connector.line.color.rgb = stroke
            connector.line.width = Pt(float(element.get("stroke_width", 2)))
            return

        mso = MSO_SHAPE.OVAL if shape_kind == "ellipse" else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(
            mso, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        if fill_hex:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _hex_to_rgb(str(fill_hex))
        else:
            shape.fill.background()
        shape.line.color.rgb = stroke
        shape.line.width = Pt(float(element.get("stroke_width", 2)))
        shape.shadow.inherit = False
