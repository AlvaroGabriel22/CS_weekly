"""Conversor de um .pptx (feito à mão no PowerPoint) para o DeckLayout interno.

Objetivo: permitir que o usuário use um weekly antigo (anterior ao sistema)
como MODELO para a IA. O modelo é um ESQUELETO — posições, fontes, cores e
onde ficam tabelas/imagens. O conteúdo real de cada semana é preenchido depois
por `_template_deck` (troca título/descrição e encaixa os anexos da semana nos
slots de tabela/imagem).

Fidelidade priorizada para o caso real dos usuários: caixas de texto, tabelas,
e imagens/gráficos colados como figura. Recursos complexos (SmartArt, gráficos
nativos não-imagem, animações) são degradados com segurança para um slot de
imagem na mesma posição — nunca quebram a importação.

Tratamento de erros: qualquer shape/slide problemático é ignorado e logado; a
importação só falha se o arquivo não for um PPTX legível ou se nada de útil for
extraído (aí devolve uma mensagem clara).
"""
from __future__ import annotations

# python-pptx 0.6.x referencia collections.abc via um caminho que quebra no
# Python 3.12 se collections.abc ainda não tiver sido importado. Garante isso.
import collections.abc  # noqa: F401
import logging
import re
from typing import Any

logger = logging.getLogger(__name__)

# Limites defensivos contra arquivos patológicos.
MAX_SLIDES = 30
MAX_ELEMENTS_PER_SLIDE = 60

DEFAULT_TEXT_COLOR = "#1F2937"
BRAND = "#0C379C"

# Vocabulário de SLOTS: o que cada elemento do modelo representa. A importação
# apenas SUGERE; quem decide é o usuário, na tela de Templates.
#
# `static` é a decoração/rótulo fixo do modelo — repete igual em todo slide e
# nunca recebe conteúdo. Todo o resto, se não for preenchido na semana, é
# limpo: é isso que impede o texto da semana anterior de vazar para o deck novo.
SLOTS = ("title", "body", "activity_date", "table", "image", "chart",
         "week_label", "static")

# "W29", "S33" — selo de semana.
_WEEK_BADGE = re.compile(r"^(w|s)\s*\d{1,2}([\s·|/-]+\d{2,4})?$", re.IGNORECASE)
# "19/07–19/07", "01/08 a 07/08", "Weekly Report - 19/07/2026 a 19/07/2026":
# período da semana, que precisa ser trocado (e não repetido) na geração.
_WEEK_RANGE = re.compile(
    r"\d{1,2}/\d{1,2}(/\d{2,4})?\s*(a|à|–|—|-|até|to)\s*\d{1,2}/\d{1,2}", re.IGNORECASE
)
# "20/07", "20/07/2026" sozinho numa caixa: a data DAQUELA atividade, que
# também precisa ser trocada (modelos costumam ter uma por bloco).
_ACTIVITY_DATE = re.compile(r"^\d{1,2}/\d{1,2}(/\d{2,4})?$")
# Um rótulo de período é uma linha curta; num parágrafo, a data é só conteúdo.
WEEK_LABEL_MAX_CHARS = 60
TITLE_MAX_CHARS = 70
BODY_MIN_CHARS = 40


def _is_week_label(text: str) -> bool:
    text = (text or "").strip()
    if not text or len(text) > WEEK_LABEL_MAX_CHARS:
        return False
    return bool(_WEEK_BADGE.match(text) or _WEEK_RANGE.search(text))


class PptxImportError(Exception):
    """Erro de importação com mensagem pronta para o usuário (PT)."""


def _frac(value: int, total: int) -> float:
    if not total:
        return 0.0
    return round(min(max(value / total, 0.0), 1.0), 4)


def _rgb_hex(font) -> str | None:
    """Cor RGB explícita da fonte (ignora cores de tema, que não resolvemos)."""
    try:
        from pptx.enum.dml import MSO_COLOR_TYPE

        color = font.color
        # Só lemos cor RGB direta; MSO_THEME_COLOR cai no default.
        if color is not None and color.type == MSO_COLOR_TYPE.RGB:
            rgb = color.rgb
            if rgb is not None:
                return f"#{str(rgb)}"
    except Exception:
        pass
    return None


def _paragraph_align(paragraph) -> str | None:
    try:
        from pptx.enum.text import PP_ALIGN

        mapping = {
            PP_ALIGN.LEFT: "left",
            PP_ALIGN.CENTER: "center",
            PP_ALIGN.RIGHT: "right",
            PP_ALIGN.JUSTIFY: "left",
        }
        return mapping.get(paragraph.alignment)
    except Exception:
        return None


def _first_run(text_frame):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            return run, paragraph
    # parágrafo sem runs (texto direto): devolve o parágrafo
    for paragraph in text_frame.paragraphs:
        return None, paragraph
    return None, None


def _text_element(shape, slide_w: int, slide_h: int, idx: int) -> dict | None:
    """Converte uma caixa de texto em elemento de texto do DeckLayout."""
    try:
        tf = shape.text_frame
    except Exception:
        return None
    text = (shape.text or "").strip()
    if not text:
        return None

    run, paragraph = _first_run(tf)

    # Tamanho da fonte (pt). Ausente = herdado do master → estimamos.
    font_size = None
    bold = False
    italic = False
    font_family = None
    color = None
    if run is not None:
        try:
            if run.font.size is not None:
                font_size = int(run.font.size.pt)
        except Exception:
            font_size = None
        bold = bool(run.font.bold)
        italic = bool(run.font.italic)
        font_family = run.font.name or None
        color = _rgb_hex(run.font)

    x = _frac(shape.left or 0, slide_w)
    y = _frac(shape.top or 0, slide_h)
    w = _frac(shape.width or 0, slide_w)
    h = _frac(shape.height or 0, slide_h)

    # Estimativa de tamanho quando o PPT herda do tema (comum). Título costuma
    # ficar no terço superior e ser curto → recebe fonte maior para que o
    # clonador o reconheça como título.
    if not font_size:
        is_titleish = y < 0.30 and len(text) <= 70
        font_size = 28 if is_titleish else 16

    element: dict[str, Any] = {
        "id": f"t{idx}",
        "type": "text",
        "x": x, "y": y,
        "w": max(w, 0.05), "h": max(h, 0.04),
        "text": text[:2000],
        "font_size": max(8, min(font_size, 60)),
        "bold": bold,
        "italic": italic,
        "color": color or DEFAULT_TEXT_COLOR,
    }
    if font_family:
        element["font_family"] = font_family
    align = _paragraph_align(paragraph) if paragraph is not None else None
    if align:
        element["align"] = align
    return element


def _slot_element(shape, slide_w: int, slide_h: int, idx: int, kind: str) -> dict:
    """Slot de tabela/imagem: guarda só posição/tamanho (o conteúdo da semana
    é encaixado depois). attachment_id fica vazio de propósito."""
    return {
        "id": f"{kind[0]}{idx}",
        "type": kind,  # "table" | "image"
        "x": _frac(shape.left or 0, slide_w),
        "y": _frac(shape.top or 0, slide_h),
        "w": max(_frac(shape.width or 0, slide_w), 0.08),
        "h": max(_frac(shape.height or 0, slide_h), 0.06),
        "attachment_id": "",
        "font_size": 12,
    }


def _shape_element(shape, slide_w: int, slide_h: int, idx: int) -> dict | None:
    """Retângulo/linha/oval decorativo (autoshape sem texto)."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            return None
    except Exception:
        return None

    w = _frac(shape.width or 0, slide_w)
    h = _frac(shape.height or 0, slide_h)
    form = "rect"
    if h <= 0.01 or w <= 0.01:
        form = "line"
    fill = None
    try:
        from pptx.enum.dml import MSO_FILL_TYPE, MSO_COLOR_TYPE

        if shape.fill.type == MSO_FILL_TYPE.SOLID:
            fore = shape.fill.fore_color
            if fore.type == MSO_COLOR_TYPE.RGB and fore.rgb is not None:
                fill = f"#{str(fore.rgb)}"
    except Exception:
        fill = None
    return {
        "id": f"s{idx}",
        "type": "shape",
        "shape": form,
        "x": _frac(shape.left or 0, slide_w),
        "y": _frac(shape.top or 0, slide_h),
        "w": max(w, 0.005), "h": max(h, 0.005),
        "color": fill or BRAND,
        "fill": fill,
        "stroke_width": 2,
        "font_size": 12,
    }


def _convert_shape(shape, slide_w: int, slide_h: int, idx: int) -> dict | None:
    """Um shape → um elemento do DeckLayout (ou None se ignorável)."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        stype = shape.shape_type
        element: dict | None
        # Tabela
        if getattr(shape, "has_table", False):
            element = _slot_element(shape, slide_w, slide_h, idx, "table")
            element["slot"] = "table"
        # Gráfico NATIVO: o exportador por mutação preenche os dados dele
        # preservando tipo, cores e estilo do modelo — por isso é slot próprio,
        # e não "mais uma imagem".
        elif getattr(shape, "has_chart", False) or stype == MSO_SHAPE_TYPE.CHART:
            element = _slot_element(shape, slide_w, slide_h, idx, "image")
            element["slot"] = "chart"
        # Imagem (inclui gráfico colado como figura)
        elif stype in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE):
            element = _slot_element(shape, slide_w, slide_h, idx, "image")
            element["slot"] = "image"
        # OLE/SmartArt → degrada para slot de imagem (posição preservada)
        elif stype == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
            element = _slot_element(shape, slide_w, slide_h, idx, "image")
            element["slot"] = "image"
        # Texto (o slot é sugerido depois, olhando o slide inteiro)
        elif getattr(shape, "has_text_frame", False) and (shape.text or "").strip():
            element = _text_element(shape, slide_w, slide_h, idx)
        # Autoshape decorativa
        else:
            element = _shape_element(shape, slide_w, slide_h, idx)
            if element:
                element["slot"] = "static"

        if element is not None:
            # Âncora para o exportador por mutação: sem isto não há como levar a
            # marcação de slot de volta ao shape certo dentro do .pptx.
            try:
                element["src_shape_id"] = int(shape.shape_id)
            except Exception:
                element["src_shape_id"] = None
        return element
    except Exception as error:
        logger.warning("Shape ignorado na importação do PPT: %s", error)
        return None


def _suggest_slots(elements: list[dict]) -> None:
    """Sugere o papel de cada caixa de TEXTO do slide (in-place).

    Regras determinísticas, e explicitamente uma SUGESTÃO — a auditoria do
    template real mostrou que "maior fonte = título" escolhe o selo da semana
    quando ele é grande. Por isso desempatamos pela LARGURA (título costuma
    atravessar o slide) e o usuário corrige na tela.
    """
    texts = [e for e in elements if e.get("type") == "text" and not e.get("slot")]
    if not texts:
        return

    for element in texts:
        texto = (element.get("text") or "").strip()
        if _is_week_label(texto):
            element["slot"] = "week_label"
        elif _ACTIVITY_DATE.match(texto):
            element["slot"] = "activity_date"

    livres = [e for e in texts if not e.get("slot")]
    if not livres:
        return

    # Só sugerimos quando o elemento é PLAUSÍVEL. Um palpite errado é pior que
    # palpite nenhum: o usuário pode não notar, e o slide sai com o conteúdo no
    # lugar trocado. Nem todo slide tem título (há moldes que são só gráfico +
    # tabela) — nesses, não marcamos nada e ele decide na tela.
    titulo = None
    candidatos = [
        e for e in livres
        if 3 < len(e.get("text") or "") <= TITLE_MAX_CHARS and (e.get("y") or 0) < 0.5
    ]
    if candidatos:
        titulo = min(candidatos, key=lambda e: (-(e.get("font_size") or 0),
                                                e.get("y") or 0,
                                                len(e.get("text") or "")))
        titulo["slot"] = "title"

    corpo = [
        e for e in livres
        if e is not titulo and len(e.get("text") or "") >= BODY_MIN_CHARS
    ]
    if corpo:
        max(corpo, key=lambda e: len(e.get("text") or ""))["slot"] = "body"

    for element in texts:
        element.setdefault("slot", "static")


def _iter_shapes(shapes):
    """Achata grupos (recursivo) para não perder elementos agrupados."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _iter_shapes(shape.shapes)
                continue
        except Exception:
            pass
        yield shape


def import_pptx_to_layout(file_path: str) -> dict:
    """Lê o .pptx e devolve um DeckLayout {slides:[...]}.

    Levanta PptxImportError (mensagem PT) se o arquivo não for legível ou se
    nenhum conteúdo aproveitável for extraído.
    """
    try:
        from pptx import Presentation
    except Exception as error:  # dependência ausente
        raise PptxImportError("Leitura de PPT indisponível no servidor.") from error

    try:
        presentation = Presentation(file_path)
    except Exception as error:
        raise PptxImportError(
            "Não foi possível abrir o arquivo. Envie um .pptx válido."
        ) from error

    slide_w = int(presentation.slide_width or 0)
    slide_h = int(presentation.slide_height or 0)
    if not slide_w or not slide_h:
        raise PptxImportError("O PPT não tem dimensões de slide válidas.")

    slides: list[dict] = []
    for s_idx, slide in enumerate(list(presentation.slides)[:MAX_SLIDES]):
        elements: list[dict] = []
        for e_idx, shape in enumerate(_iter_shapes(slide.shapes)):
            if len(elements) >= MAX_ELEMENTS_PER_SLIDE:
                break
            element = _convert_shape(shape, slide_w, slide_h, e_idx)
            if element:
                elements.append(element)
        if not elements:
            continue  # slide vazio/ilegível — pula
        _suggest_slots(elements)
        for element in elements:
            element["src_slide"] = s_idx
        slides.append({
            "id": f"s{s_idx}",
            # 1º slide com conteúdo = capa; demais = molde de conteúdo.
            "kind": "cover" if not slides else "custom",
            "elements": elements,
            "src_slide": s_idx,
        })

    if not slides:
        raise PptxImportError(
            "Não encontramos texto, tabelas ou imagens neste PPT. "
            "Verifique se o arquivo tem conteúdo em caixas de texto."
        )
    return {"slides": slides}
