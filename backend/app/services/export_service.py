"""Export Service - PPTX report generation"""
import os
from datetime import datetime
from typing import List, Optional
from sqlalchemy.orm import Session
from app.models import WeeklyReport, Activity
from app.repositories import ActivityRepository


class ExportService:
    """Service for exporting reports to PPTX"""

    EXPORT_DIR = 'backend/uploads/reports'

    def __init__(self, db: Session):
        self.db = db
        self.activity_repo = ActivityRepository(db)
        os.makedirs(self.EXPORT_DIR, exist_ok=True)

    def generate_pptx(
        self,
        report: WeeklyReport,
        activities: List[Activity],
        title: Optional[str] = None,
    ) -> str:
        """Generate PPTX presentation from report"""

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
        except ImportError:
            raise ImportError("python-pptx is required for PPTX export. Install with: pip install python-pptx")

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Title slide
        self._add_title_slide(prs, report, title)

        # Summary slide
        self._add_summary_slide(prs, report, activities)

        # Activities slides
        for activity in activities:
            self._add_activity_slide(prs, activity)

        # Statistics slide
        self._add_statistics_slide(prs, report, activities)

        # Save presentation
        filename = f"relatorio_w{report.week_number}_{report.year}.pptx"
        pptx_path = os.path.join(self.EXPORT_DIR, filename)

        prs.save(pptx_path)
        return pptx_path

    @staticmethod
    def _add_title_slide(prs, report: WeeklyReport, title: Optional[str]) -> None:
        """Add title slide to presentation"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(100, 150, 200)  # Blue pastel

        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title or f"Relatório Semanal - Semana {report.week_number}/{report.year}"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Date
        date_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
        date_frame = date_box.text_frame
        p = date_frame.paragraphs[0]
        p.text = f"Gerado em {datetime.now().strftime('%d de %B de %Y')}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    @staticmethod
    def _add_summary_slide(prs, report: WeeklyReport, activities: List[Activity]) -> None:
        """Add summary slide"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
        title = slide.shapes.title
        title.text = "Resumo da Semana"

        # Summary content
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        # Statistics
        summary_text = f"""
Período: Semana {report.week_number}/{report.year}

Total de Atividades: {len(activities)}
Atividades Registradas: {sum(1 for a in activities if a.status.value == 'registered')}
Atividades Processadas: {sum(1 for a in activities if a.status.value == 'processed')}

Arquivos Anexados: {sum(len(a.attachments) for a in activities if a.attachments)}
Atividades com Metadados: {sum(1 for a in activities if a.metadata_entry)}

Status: Gerado automaticamente pelo QWI
        """

        p = tf.paragraphs[0]
        p.text = summary_text.strip()
        p.font.size = Pt(14)
        p.line_spacing = 1.5

    @staticmethod
    def _add_activity_slide(prs, activity: Activity) -> None:
        """Add slide for single activity"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = activity.title[:50]  # Truncate long titles

        # Activity details
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        # Content
        content = f"""
{activity.description or 'Sem descrição'}

Status: {activity.status.value}
Tags: {', '.join(activity.tags) if activity.tags else 'Nenhuma'}
Projeto: {activity.project or 'N/A'}
Categoria: {activity.category or 'N/A'}

Arquivos: {len(activity.attachments) if activity.attachments else 0}
Metadados: {'Sim' if activity.metadata_entry else 'Não'}
        """

        p = tf.paragraphs[0]
        p.text = content.strip()
        p.font.size = Pt(12)
        p.line_spacing = 1.4

    @staticmethod
    def _add_statistics_slide(prs, report: WeeklyReport, activities: List[Activity]) -> None:
        """Add statistics slide"""
        from pptx.util import Inches, Pt

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Estatísticas"

        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        # Calculate stats
        by_status = {}
        for activity in activities:
            status = activity.status.value
            by_status[status] = by_status.get(status, 0) + 1

        stats_text = f"""
Por Status:
"""
        for status, count in by_status.items():
            stats_text += f"  • {status}: {count}\n"

        stats_text += f"""
Arquivos Totais: {sum(len(a.attachments) for a in activities if a.attachments)}
Atividades com Metadados: {sum(1 for a in activities if a.metadata_entry)}

Qualidade do Relatório: {report.quality_score * 100 if report.quality_score else 'N/A'}%
Confiança: {'Alta' if report.confidence_index and max(report.confidence_index) > 0.8 else 'Média'}
        """

        p = tf.paragraphs[0]
        p.text = stats_text.strip()
        p.font.size = Pt(12)
