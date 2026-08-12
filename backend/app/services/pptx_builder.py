"""Comprehensive PPTX builder with template support and advanced customization."""

import io
import logging
import os
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any, Optional

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from app.core.config import get_settings
from app.services.chart_service import ChartService
from app.services.pptx_templates import TemplateEngine, TemplateTheme

logger = logging.getLogger(__name__)
settings = get_settings()


@dataclass
class SlideConfig:
    """Configuration for a slide."""
    title: str
    subtitle: Optional[str] = None
    layout: str = "default"  # "title", "content", "two_column", "blank"
    background_color: Optional[tuple] = None  # RGB tuple
    theme: Optional[str] = None


class PPTXBuilder:
    """High-level PPTX generator with template support and customization."""

    def __init__(self, template_theme: Optional[TemplateTheme] = None, output_dir: Optional[str] = None):
        """Initialize the PPTX builder.

        Args:
            template_theme: Optional theme configuration
            output_dir: Output directory for generated PPTX files
        """
        self.template_engine = TemplateEngine()
        self.chart_service = ChartService()
        self.theme = template_theme or TemplateTheme.default()

        upload_root = Path(settings.UPLOAD_DIR)
        self.output_dir = Path(output_dir) if output_dir else upload_root / "reports"
        self.output_dir.mkdir(parents=True, exist_ok=True)

        self.presentation: Optional[Presentation] = None
        self.slide_count = 0

    def create_presentation(self, width_inches: float = 10, height_inches: float = 7.5) -> Presentation:
        """Create a new presentation with custom dimensions."""
        self.presentation = Presentation()
        self.presentation.slide_width = Inches(width_inches)
        self.presentation.slide_height = Inches(height_inches)
        self.slide_count = 0
        return self.presentation

    def generate_pptx(
        self,
        weekly_report: dict[str, Any],
        template_name: str = "default",
        filename: Optional[str] = None,
    ) -> str:
        """Generate complete PPTX from weekly report.

        Args:
            weekly_report: Weekly report data
            template_name: Template to use
            filename: Output filename (auto-generated if not provided)

        Returns:
            Path to generated PPTX file
        """
        self.create_presentation()

        # Add cover slide
        self.add_title_slide(
            title=weekly_report.get("title", "Weekly Report"),
            subtitle=weekly_report.get("department", ""),
            date=weekly_report.get("date", datetime.now().strftime("%Y-%m-%d")),
            author=weekly_report.get("author", "QWI"),
        )

        # Add summary if available
        if weekly_report.get("summary"):
            self.add_content_slide(
                title="Executive Summary",
                content=weekly_report["summary"],
                images=None,
            )

        # Add activities
        activities = weekly_report.get("activities", [])
        for idx, activity in enumerate(activities, 1):
            self._add_activity_slide(activity, idx)

        # Add KPI summary if available
        kpi_table = weekly_report.get("kpi_table")
        if kpi_table:
            self.add_table_slide(
                title="Key Performance Indicators",
                data=self._format_kpi_data(kpi_table),
            )

        # Add conclusions and next steps
        if weekly_report.get("conclusions"):
            self.add_content_slide(
                title="Conclusions",
                content="\n".join(weekly_report["conclusions"]),
            )

        if weekly_report.get("next_steps"):
            self.add_content_slide(
                title="Next Steps",
                content="\n".join(weekly_report["next_steps"]),
            )

        # Save presentation
        output_path = self._save_presentation(filename)
        logger.info("PPTX generated | path=%s | slides=%d", output_path, self.slide_count)
        return output_path

    def add_title_slide(
        self,
        title: str,
        subtitle: str = "",
        date: str = "",
        author: str = "",
    ) -> None:
        """Add a title/cover slide."""
        if not self.presentation:
            self.create_presentation()

        slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[6]  # Blank layout
        )
        self._set_background(slide, self.theme.background_color)

        # Add accent bar
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            Inches(0.25),
            self.presentation.slide_height,
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(*self.theme.accent_color)
        accent.line.fill.background()

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.9), Inches(2.0),
            Inches(8.5), Inches(1.5),
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*self.theme.text_color)

        # Add subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.9), Inches(3.6),
                Inches(8.5), Inches(0.4),
            )
            subtitle_frame = subtitle_box.text_frame
            p = subtitle_frame.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(*self.theme.secondary_text_color)

        # Add date
        if date:
            date_box = slide.shapes.add_textbox(
                Inches(0.9), Inches(4.1),
                Inches(8.5), Inches(0.3),
            )
            date_frame = date_box.text_frame
            p = date_frame.paragraphs[0]
            p.text = date
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(*self.theme.secondary_text_color)

        # Add author
        if author:
            author_box = slide.shapes.add_textbox(
                Inches(0.9), Inches(6.0),
                Inches(8.5), Inches(0.3),
            )
            author_frame = author_box.text_frame
            p = author_frame.paragraphs[0]
            p.text = f"Author: {author}"
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(*self.theme.secondary_text_color)

        self.slide_count += 1

    def add_content_slide(
        self,
        title: str,
        content: str,
        images: Optional[list[str]] = None,
    ) -> None:
        """Add a content slide with title and text."""
        if not self.presentation:
            self.create_presentation()

        slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[6]  # Blank layout
        )
        self._set_background(slide, self.theme.background_color)

        # Add header with accent line
        self._add_slide_header(slide, title)

        # Add content text
        content_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.2),
            Inches(8.5), Inches(5.0),
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        # Parse content (support bullet points with "- " prefix)
        lines = content.split("\n")
        for i, line in enumerate(lines):
            if i > 0:
                content_frame.add_paragraph()
            p = content_frame.paragraphs[i]
            p.text = line.lstrip("- ").strip()
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(*self.theme.text_color)
            if line.strip().startswith("-"):
                p.level = 0

        # Add images if provided
        if images:
            self._add_slide_images(slide, images, Inches(6.5), Inches(1.2))

        self.slide_count += 1

    def add_chart_slide(
        self,
        title: str,
        data: dict[str, Any],
        chart_type: str = "bar",
    ) -> None:
        """Add a slide with a chart.

        Args:
            title: Slide title
            data: Chart data (categories, series, values)
            chart_type: Chart type (bar, pie, line, column)
        """
        if not self.presentation:
            self.create_presentation()

        slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[6]  # Blank layout
        )
        self._set_background(slide, self.theme.background_color)

        # Add header
        self._add_slide_header(slide, title)

        # Generate chart image
        chart_image_path = self.chart_service.generate_chart(
            data=data,
            chart_type=chart_type,
            theme=self.theme,
        )

        if chart_image_path and Path(chart_image_path).exists():
            # Add chart image to slide
            slide.shapes.add_picture(
                chart_image_path,
                Inches(0.75), Inches(1.2),
                width=Inches(8.5),
            )

        self.slide_count += 1

    def add_table_slide(
        self,
        title: str,
        data: list[list[str]],
        headers: Optional[list[str]] = None,
    ) -> None:
        """Add a slide with a table.

        Args:
            title: Slide title
            data: Table data (list of rows)
            headers: Optional column headers
        """
        if not self.presentation:
            self.create_presentation()

        slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[6]  # Blank layout
        )
        self._set_background(slide, self.theme.background_color)

        # Add header
        self._add_slide_header(slide, title)

        # Calculate table dimensions
        rows = len(data) + (1 if headers else 0)
        cols = len(data[0]) if data else 1

        # Add table
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(9)
        height = Inches(4.5)

        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set column widths
        col_width = width / cols
        for col in range(cols):
            table_shape.columns[col].width = int(col_width)

        # Add headers if provided
        if headers:
            for col_idx, header_text in enumerate(headers):
                cell = table_shape.cell(0, col_idx)
                cell.text = header_text
                self._format_table_cell(cell, bold=True, bg_color=self.theme.accent_color)

        # Add data rows
        start_row = 1 if headers else 0
        for row_idx, row_data in enumerate(data, start=start_row):
            for col_idx, cell_text in enumerate(row_data):
                cell = table_shape.cell(row_idx, col_idx)
                cell.text = str(cell_text)
                self._format_table_cell(cell)

        self.slide_count += 1

    def add_image_slide(
        self,
        title: str,
        images: list[str],
    ) -> None:
        """Add a slide with multiple images.

        Args:
            title: Slide title
            images: List of image paths
        """
        if not self.presentation:
            self.create_presentation()

        slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[6]  # Blank layout
        )
        self._set_background(slide, self.theme.background_color)

        # Add header
        self._add_slide_header(slide, title)

        # Add images in grid
        self._add_image_grid(slide, images, Inches(0.75), Inches(1.2))

        self.slide_count += 1

    def add_summary_slide(
        self,
        activities_summary: dict[str, Any],
    ) -> None:
        """Add a summary slide with activity statistics.

        Args:
            activities_summary: Summary statistics (counts, statuses, etc.)
        """
        if not self.presentation:
            self.create_presentation()

        slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[6]  # Blank layout
        )
        self._set_background(slide, self.theme.background_color)

        # Add header
        self._add_slide_header(slide, "Summary Statistics")

        # Add summary data
        y_pos = Inches(1.2)
        for key, value in activities_summary.items():
            # Key-value pair
            text_box = slide.shapes.add_textbox(
                Inches(0.75), y_pos,
                Inches(8.5), Inches(0.3),
            )
            text_frame = text_box.text_frame
            p = text_frame.paragraphs[0]
            p.text = f"{key}: {value}"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(*self.theme.text_color)
            y_pos += Inches(0.4)

        self.slide_count += 1

    # ==================== Private helper methods ====================

    def _add_activity_slide(self, activity: dict[str, Any], index: int) -> None:
        """Add an activity slide with all its details."""
        slide = self.presentation.slides.add_slide(
            self.presentation.slide_layouts[6]  # Blank layout
        )
        self._set_background(slide, self.theme.background_color)

        title = f"Activity {index}: {activity.get('title', 'Untitled')}"
        self._add_slide_header(slide, title)

        # Add narrative
        narrative = activity.get("narrative", "")
        if narrative:
            content_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(1.2),
                Inches(5.5), Inches(5.0),
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            p = content_frame.paragraphs[0]
            p.text = narrative
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(*self.theme.text_color)

        # Add images if available
        images = activity.get("images", [])
        if images:
            image_paths = [img.get("path") or img.get("url") for img in images if img.get("path") or img.get("url")]
            if image_paths:
                self._add_slide_images(slide, image_paths, Inches(6.3), Inches(1.2))

        # Add impact if available
        impact = activity.get("impact", "")
        if impact:
            impact_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(5.8),
                Inches(5.5), Inches(1.0),
            )
            impact_frame = impact_box.text_frame
            impact_frame.word_wrap = True
            p = impact_frame.paragraphs[0]
            p.text = f"Impact: {impact}"
            p.font.size = Pt(10)
            p.font.italic = True
            p.font.color.rgb = RGBColor(*self.theme.secondary_text_color)

        self.slide_count += 1

    def _add_slide_header(self, slide, title: str) -> None:
        """Add a header with title and accent bar."""
        # Add horizontal line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.75), Inches(0.95),
            Inches(8.5), Inches(0.03),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(*self.theme.accent_color)
        line.line.fill.background()

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(0.5),
            Inches(8.5), Inches(0.4),
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*self.theme.text_color)

    def _add_slide_images(
        self,
        slide,
        images: list[str],
        left: Inches,
        top: Inches,
        max_width: Inches = Inches(3),
    ) -> None:
        """Add images to a slide with auto-resize."""
        current_y = top
        for image_path in images:
            if not Path(image_path).exists():
                logger.warning("Image not found: %s", image_path)
                continue

            try:
                # Get image dimensions
                img = Image.open(image_path)
                img_width, img_height = img.size

                # Calculate scaled size
                aspect_ratio = img_height / img_width
                scaled_width = max_width
                scaled_height = scaled_width * aspect_ratio

                # Check if image fits vertically
                if current_y + scaled_height > Inches(6.5):
                    break

                # Add image
                slide.shapes.add_picture(
                    image_path,
                    left, current_y,
                    width=scaled_width,
                )
                current_y += scaled_height + Inches(0.2)
            except Exception as e:
                logger.error("Failed to add image %s: %s", image_path, e)

    def _add_image_grid(
        self,
        slide,
        images: list[str],
        left: Inches,
        top: Inches,
        cols: int = 2,
        img_width: Inches = Inches(4),
    ) -> None:
        """Add images in a grid layout."""
        current_y = top
        current_x = left
        col_count = 0

        for image_path in images:
            if not Path(image_path).exists():
                logger.warning("Image not found: %s", image_path)
                continue

            if col_count >= cols:
                current_y += img_width * 0.75 + Inches(0.2)
                current_x = left
                col_count = 0

            try:
                slide.shapes.add_picture(
                    image_path,
                    current_x, current_y,
                    width=img_width,
                )
                current_x += img_width + Inches(0.2)
                col_count += 1
            except Exception as e:
                logger.error("Failed to add image %s: %s", image_path, e)

    def _format_table_cell(
        self,
        cell,
        text: Optional[str] = None,
        bold: bool = False,
        bg_color: Optional[tuple] = None,
    ) -> None:
        """Format a table cell."""
        if text:
            cell.text = text

        if bg_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*bg_color)

        # Format text
        if cell.text_frame.paragraphs:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    if bold:
                        run.font.bold = True
                    run.font.color.rgb = RGBColor(*self.theme.text_color)

    def _set_background(self, slide, color: Optional[tuple] = None) -> None:
        """Set slide background color."""
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(*(color or self.theme.background_color))

    def _format_kpi_data(self, kpi_table: list[dict]) -> list[list[str]]:
        """Format KPI table data for display."""
        data = []
        for row in kpi_table:
            data.append([
                row.get("kpi", ""),
                row.get("result", ""),
                row.get("trend", ""),
            ])
        return data

    def _save_presentation(self, filename: Optional[str] = None) -> str:
        """Save the presentation and return its path."""
        if not self.presentation:
            raise ValueError("No presentation created")

        if not filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"weekly_report_{timestamp}.pptx"

        output_path = self.output_dir / filename
        self.presentation.save(str(output_path))
        return str(output_path)

    def export_to_pdf(self, pptx_path: str, pdf_path: Optional[str] = None) -> Optional[str]:
        """Export PPTX to PDF (requires LibreOffice).

        Note: This is a stub that requires external tool integration.
        """
        try:
            import subprocess

            if not pdf_path:
                pdf_path = str(Path(pptx_path).with_suffix(".pdf"))

            # Use LibreOffice to convert
            subprocess.run([
                "libreoffice", "--headless", "--convert-to", "pdf",
                "--outdir", str(Path(pdf_path).parent),
                pptx_path,
            ], check=True)

            return pdf_path
        except Exception as e:
            logger.error("PDF export failed: %s", e)
            return None

    def export_to_docx(self, pptx_path: str) -> Optional[str]:
        """Export PPTX to DOCX (stub for future implementation)."""
        logger.warning("DOCX export not yet implemented")
        return None
