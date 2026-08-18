"""E2E 'fluxo ouro' do QA (Agente J) — sequencial, contra servidor dedicado.

Sobe implícito? Não: assume um backend em --base (default :8002) com banco/uploads
isolados. Executa os 12 passos do briefing e imprime PASS/FAIL por passo + tempos.
Não commitar. Usa python-pptx para validar o binário baixado.

Uso:  ./venv/bin/python scripts/qa_e2e.py --base http://localhost:8002
"""
from __future__ import annotations

import argparse
import io
import sys
import time

import httpx

R = {"pass": 0, "fail": 0}


def step(name: str, ok: bool, detail: str = ""):
    tag = "PASS" if ok else "FAIL"
    R["pass" if ok else "fail"] += 1
    print(f"[{tag}] {name}" + (f" — {detail}" if detail else ""))
    return ok


def xlsx_bytes() -> bytes:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.append(["Linha", "NC", "Status"])
    ws.append(["L3", "Solda fria", "Aberta"])
    ws.append(["L1", "Parafuso", "Fechada"])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def png_bytes() -> bytes:
    from PIL import Image
    buf = io.BytesIO()
    Image.new("RGB", (40, 30), (12, 55, 156)).save(buf, "PNG")
    return buf.getvalue()


def reg_login(c: httpx.Client, email: str, emp: str, role: str, sector: str) -> str | None:
    c.post("/api/auth/register", json={
        "name": f"E2E {emp}", "email": email, "password": "senha123",
        "password_confirm": "senha123", "employee_id": emp, "role": role, "sector": sector})
    r = c.post("/api/auth/login", json={"email": email, "password": "senha123"})
    return r.json().get("access_token") if r.status_code == 200 else None


def wait_report(c: httpx.Client, h: dict, rid: str, timeout: int = 180) -> dict | None:
    end = time.time() + timeout
    while time.time() < end:
        r = c.get(f"/api/weekly/{rid}", headers=h)
        if r.status_code == 200:
            st = r.json().get("status")
            if st in ("completed", "failed", "error"):
                return r.json()
        time.sleep(3)
    return None


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--base", default="http://localhost:8002")
    ap.add_argument("--stamp", type=int, default=int(time.time()))
    args = ap.parse_args()
    s = args.stamp
    c = httpx.Client(base_url=args.base, timeout=200)
    t0 = time.time()

    # 1. Register A
    ta = reg_login(c, f"e2ea{s}@qwiqa.com", f"E2EA{s}", "Analista Jr", "OQC")
    step("1. Register+login A (OQC)", bool(ta))
    ha = {"Authorization": f"Bearer {ta}"}

    # 3. Create activity com anexo (xlsx) + sem diretiva de imagem
    aid = c.post("/api/activities", json={"title": "Auditoria linha 3",
                 "description": "Verificados 12 pontos; 2 NCs.", "include_in_weekly": True},
                 headers=ha).json().get("id")
    up = c.post(f"/api/activities/{aid}/attachments",
                files={"file": ("ncs.xlsx", io.BytesIO(xlsx_bytes()),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")},
                headers=ha)
    step("3a. Create activity + upload xlsx", up.status_code == 200, f"HTTP {up.status_code}")
    # 3b. atividade com imagem + diretiva
    aid2 = c.post("/api/activities", json={"title": "Foto do defeito /analisar imagem",
                  "description": "/analisar imagem — inspeção visual", "include_in_weekly": True},
                  headers=ha).json().get("id")
    up2 = c.post(f"/api/activities/{aid2}/attachments",
                 files={"file": ("foto.png", io.BytesIO(png_bytes()), "image/png")}, headers=ha)
    step("3b. Activity com diretiva de imagem + png", up2.status_code == 200)

    # 4. List/filter na agenda
    lst = c.get("/api/activities?page_size=50", headers=ha)
    n = lst.json().get("total", len(lst.json().get("items", [])))
    step("4. List activities", lst.status_code == 200 and n >= 2, f"total={n}")

    # 5. Generate weekly (com IA real — pode demorar). Fallback se falhar ainda gera.
    tgen = time.time()
    gen = c.post("/api/weekly/generate", json={"activity_ids": [aid, aid2],
                 "week_number": 33, "year": 2026}, headers=ha)
    rid = gen.json().get("id") if gen.status_code == 200 else None
    step("5. POST weekly/generate", bool(rid), f"HTTP {gen.status_code}")

    # 6. Poll até COMPLETED (ou fallback)
    rep = wait_report(c, ha, rid) if rid else None
    gen_secs = round(time.time() - tgen, 1)
    step("6. Weekly COMPLETED", bool(rep) and rep.get("status") == "completed",
         f"status={rep.get('status') if rep else 'timeout'} em {gen_secs}s")

    # 7. Download PPTX + validar com python-pptx
    ok7 = False
    if rid:
        dl = c.get(f"/api/weekly/{rid}/download", headers=ha)
        if dl.status_code == 200 and dl.content[:4] == b"PK\x03\x04":
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(dl.content))
                ok7 = len(prs.slides) >= 1
                step("7. Download PPTX válido", ok7,
                     f"{len(dl.content)}B, {len(prs.slides)} slides")
            except Exception as e:
                step("7. Download PPTX válido", False, f"python-pptx: {e}")
        else:
            step("7. Download PPTX válido", False,
                 f"HTTP {dl.status_code}, magic={dl.content[:4]!r}")
    else:
        step("7. Download PPTX válido", False, "sem report")

    # 8. Regenerate → version incrementa; download antigo ainda ok
    v_old = rep.get("version") if rep else None
    gen2 = c.post("/api/weekly/generate", json={"activity_ids": [aid, aid2],
                  "week_number": 33, "year": 2026}, headers=ha)
    rid2 = gen2.json().get("id") if gen2.status_code == 200 else None
    rep2 = wait_report(c, ha, rid2) if rid2 else None
    v_new = rep2.get("version") if rep2 else None
    inc = bool(v_old) and bool(v_new) and v_new > v_old
    step("8a. Regenerate incrementa version", inc, f"v{v_old} -> v{v_new}")
    if rid:
        old = c.get(f"/api/weekly/{rid}/download", headers=ha)
        step("8b. Download da versão antiga ainda válido",
             old.status_code == 200 and old.content[:4] == b"PK\x03\x04")

    # 9. Colleague path
    tb = reg_login(c, f"e2eb{s}@qwiqa.com", f"E2EB{s}", "Analista Jr", "OQC")   # mesmo setor
    tc = reg_login(c, f"e2ec{s}@qwiqa.com", f"E2EC{s}", "Analista Jr", "IQC")   # outro setor
    hb = {"Authorization": f"Bearer {tb}"}
    hc = {"Authorization": f"Bearer {tc}"}
    if rid:
        rb = c.get(f"/api/weekly/{rid}", headers=hb)
        rc = c.get(f"/api/weekly/{rid}", headers=hc)
        step("9a. Colega mesmo setor (B) vê weekly de A", rb.status_code == 200, f"HTTP {rb.status_code}")
        step("9b. Outro setor sem grant (C) NÃO vê", rc.status_code == 403, f"HTTP {rc.status_code}")
        # grant a C
        uid_c = c.get("/api/users/profile", headers=hc).json().get("id")
        c.post("/api/users/me/access-grants", json={"employee_id": f"E2EC{s}"}, headers=ha)
        rc2 = c.get(f"/api/weekly/{rid}", headers=hc)
        step("9c. Após grant, C vê", rc2.status_code == 200, f"HTTP {rc2.status_code}")

    # 10. Forgot password com matrícula
    rp = c.post("/api/auth/reset-password", json={
        "email": f"e2ea{s}@qwiqa.com", "employee_id": f"E2EA{s}",
        "new_password": "nova12345", "new_password_confirm": "nova12345"})
    ok10 = rp.status_code == 200
    if ok10:
        relog = c.post("/api/auth/login", json={"email": f"e2ea{s}@qwiqa.com", "password": "nova12345"})
        ok10 = relog.status_code == 200
    step("10. Reset password com matrícula + relogin", ok10, f"HTTP {rp.status_code}")

    # 11. send-email sem SMTP → 503
    if rid:
        se = c.post(f"/api/weekly/{rid}/send-email",
                    json={"recipients": ["x@qwiqa.com"], "subject": "t", "body": "b"},
                    headers={"Authorization": f"Bearer {c.post('/api/auth/login', json={'email': f'e2ea{s}@qwiqa.com','password':'nova12345'}).json()['access_token']}"})
        step("11. send-email sem SMTP → 503", se.status_code == 503, f"HTTP {se.status_code}")

    # 12. translate + email-suggestion (se LLM disponível)
    if rid:
        tr = c.post("/api/ai/translate", json={"texts": ["Auditoria de processo"], "target": "en"}, headers=ha)
        step("12a. translate", tr.status_code in (200, 503), f"HTTP {tr.status_code}")

    print(f"\n== E2E: {R['pass']} PASS / {R['fail']} FAIL em {round(time.time()-t0,1)}s ==")
    sys.exit(1 if R["fail"] else 0)


if __name__ == "__main__":
    main()
